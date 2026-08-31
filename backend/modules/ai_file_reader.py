"""
modules/ai_file_reader.py

Dua jalur berbeda, TIDAK satu mekanisme tunggal (sama seperti cara Claude
sendiri "membaca" file):

  A. File teks/gambar/PDF kecil -- dikirim LANGSUNG ke Claude API (base64
     untuk gambar/PDF, teks polos untuk file teks), Claude yang "membaca"
     isinya sendiri, tanpa parsing manual di sisi kita.
  B. [BARU -- POIN 2] File besar/biner (xlsx/xlsm/xls, docx/doc) & dataset
     besar (CSV berukuran besar) -- TIDAK dikirim mentah ke Claude API
     (Claude API tidak menerima xlsx/docx sebagai content block apa pun,
     cuma teks/gambar/PDF), jadi diparsing DULU secara lokal pakai
     pandas/openpyxl/python-docx (xlsx/docx modern) atau xlrd/textract
     (xls/doc format lama, lihat _EKSTENSI_OFFICE_LAMA) -- analog "sandbox
     Linux" yang dipakai Claude sendiri utk file besar/biner. Hasil
     ekstraksi (ringkasan sheet/tabel/paragraf, BUKAN dump seluruh isi)
     baru dikirim sebagai teks ke Claude untuk diinterpretasi, dan
     di-CACHE di disk per hash isi file (POIN 4: "gak proses ulang yang
     gak perlu") -- lihat _FOLDER_CACHE_EKSTRAKSI_OFFICE.

[DIPERBAIKI -- lihat CATATAN PERBAIKAN di bawah] Versi ini menutup gap dari
versi dasar sebelumnya: model name usang, tidak ada validasi limit ukuran
(PDF 32MB/100 halaman, gambar 5MB/20 per request -- limit resmi Claude API),
tidak ada retry utk rate-limit/overload, cuma bisa 1 file per panggilan, dan
tidak ada dukungan xlsx/docx/dataset besar sama sekali.

Perlu: pip install anthropic
Perlu (utk .docx): pip install python-docx
Perlu: environment variable ANTHROPIC_API_KEY
"""
from __future__ import annotations

import atexit
import base64
import concurrent.futures
import hashlib
import io
import os
import pickle
import random
import re
import threading
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import anthropic
import openpyxl
import pandas as pd

from .logging_config import get_module_logger
# [GABUNG -- semua panggilan Claude sekarang lewat modules/claude_client.py]
# _ambil_client & _panggil_dengan_retry di bawah adalah ALIAS ke fungsi di
# claude_client.py (bukan didefinisikan ulang di sini) -- supaya SEMUA
# call site yang sudah ada di file ini (kirim_pdf_ke_ai, kirim_gambar_ke_ai,
# kirim_file_ke_ai, kirim_file_ke_ai_stream, kirim_banyak_file_ke_ai, dst)
# TIDAK perlu diubah satu-satu, cukup titik definisinya yang pindah.
from .claude_client import ambil_client as _ambil_client
from .claude_client import panggil_dengan_retry as _panggil_dengan_retry

logger = get_module_logger("ai_file_reader")

# [FIX -- GAP 1: MODEL NAME USANG] "claude-sonnet-4-6" sudah tidak berlaku.
# Model API terkini (per Agustus 2026): claude-sonnet-5, claude-opus-4-8,
# claude-haiku-4-5-20251001, claude-fable-5. Default dipilih Sonnet karena
# paling seimbang biaya/kualitas untuk baca dokumen -- ganti lewat parameter
# `model` di tiap fungsi kalau perlu Opus (akurasi lebih tinggi) atau Haiku
# (lebih murah/cepat untuk file sederhana).
MODEL_DEFAULT = "claude-sonnet-5"

# [FIX -- GAP 2: TIDAK ADA VALIDASI LIMIT UKURAN] Limit resmi Claude API
# (Messages API, direct/first-party) -- lihat docs.claude.com/pdf-support
# dan docs image support. Payload PDF+gambar lain dalam 1 request dihitung
# GABUNGAN ke limit request 32MB, jadi validasi di sini per-file, TAPI
# caller (endpoint FastAPI) tetap harus jaga total payload kalau kirim
# banyak file sekaligus lewat kirim_banyak_file_ke_ai().
MAX_PDF_BYTES = 32 * 1024 * 1024      # 32 MB per PDF
# [FIX -- LIMIT USANG] Sebelumnya 100 halaman, berdasarkan limit API lama
# untuk model context-window 200rb token. MODEL_DEFAULT (claude-sonnet-5)
# punya context window 1 JUTA token (default, bukan beta) -- limit resmi
# Claude API terkini untuk model context-window 1M adalah 600 halaman
# PDF/gambar per request (100 tetap berlaku HANYA untuk model 200rb token,
# mis. kalau `model` diganti ke Haiku). Nilai lama (100) sebelumnya
# menolak file yang sebenarnya masih valid dikirim ke Sonnet 5 (mis.
# rekening koran >100 halaman).
MAX_PDF_HALAMAN = 600                  # analisis visual penuh; >600 hlm ditolak di sini
                                        # (limit resmi Claude API utk model
                                        # context-window 1M token per Agustus 2026)
MAX_GAMBAR_BYTES = 5 * 1024 * 1024    # 5 MB per gambar
MAX_GAMBAR_PER_REQUEST = 20            # maksimum gambar per 1 panggilan API
MAX_KARAKTER_TEKS = 400_000            # ~100rb token kasar -- batas AMAN sebelum
                                        # dipotong, supaya tidak mendorong context
                                        # window habis sendirian (lihat catatan
                                        # di kirim_file_ke_ai untuk file teks) --
                                        # tetap longgar dibanding context window 1M
                                        # token Sonnet 5, sengaja konservatif krn ini
                                        # cuma utk 1 file teks, bukan seluruh prompt.

# [FIX -- LIMIT USANG] Sebelumnya HARDCODE max_tokens=1024 di 8 tempat
# berbeda di file ini -- warisan dari limit output lama. MODEL_DEFAULT
# (claude-sonnet-5) mendukung output HINGGA 128.000 token per request.
# 1024 token (~700-800 kata) SANGAT gampang membuat jawaban terpotong utk
# kasus pemakaian file ini sendiri (lihat blok __main__ di bawah: minta
# ekstrak SEMUA transaksi dari rekening koran jadi tabel -- rekening koran
# 25 halaman dgn ratusan transaksi jelas butuh jauh lebih dari 1024 token
# output). 8192 dipilih sbg default baru yang jauh lebih aman tanpa
# langsung memaksimalkan ke 128rb (yang menaikkan risiko biaya per request
# kalau tidak sengaja) -- naikkan lagi manual di sini kalau kasus
# pemakaian tertentu (mis. ekstraksi tabel sangat panjang) masih terpotong.
MAX_TOKENS_JAWABAN_DEFAULT = 8192

# [BARU] System prompt utk fitur AI File Reader (endpoint /api/ai-baca-file,
# /api/ai-baca-file-stream, /api/ai-baca-banyak-file -- SEMUA titik yang
# user aplikasi upload file bebas + pertanyaan bebas). SEBELUMNYA tidak ada
# system prompt sama sekali di jalur ini -- Claude cuma menerima isi file +
# pertanyaan user, tanpa tahu ini aplikasi akuntansi utk siapa & harus
# menangani jenis dokumen apa. Ini beda dari system prompt lain di
# claude_client.py (_SYSTEM_PROMPT_NARASI_CALK dkk) yang scope-nya sempit
# utk 1 tugas terstruktur internal -- system prompt ini scope-nya LEBAR
# krn endpoint ini terima SEMBARANG file dari user aplikasi.
SYSTEM_PROMPT_FILE_READER = (
    "Kamu adalah asisten AI di aplikasi akuntansi milik Gouf Consulting, "
    "dipakai akuntan untuk membaca & menganalisis dokumen klien (rekening "
    "koran/mutasi bank, data penjualan & invoice, bukti potong PPh "
    "21/23/4(2), slip gaji, data aset tetap, buku bantu piutang/hutang, "
    "rekap penilaian kinerja, dan dokumen akuntansi/perpajakan Indonesia "
    "lainnya).\n\n"
    "Sebelum menjawab, kenali dulu jenis dokumennya dari STRUKTUR isinya "
    "(kolom, header, pola baris) -- bukan dari nama file, karena tiap "
    "klien punya nama file & isi yang beda-beda. Kalau dokumennya rekening "
    "koran/mutasi bank, perhatikan pola berikut (umum di rekening koran "
    "bank Indonesia, bukan spesifik 1 bank saja): deskripsi transaksi "
    "sering berupa blok multi-baris (baris pertama = jenis channel "
    "transaksi, baris berikutnya = nomor referensi & nama pihak lawan "
    "transaksi); kolom saldo TIDAK selalu terisi di tiap baris -- baris "
    "kosong berarti saldo belum berubah dari checkpoint terakhir, BUKAN "
    "data hilang; penanda arah mutasi (debit/kredit) bisa berupa suffix "
    "teks, tanda minus, atau kolom terpisah, tergantung bank; dan biasanya "
    "ada baris ringkasan penutup (saldo awal, total mutasi kredit/debit, "
    "saldo akhir) yang bisa dipakai memvalidasi hasil ekstraksi.\n\n"
    "Kalau pertanyaan dari user terlalu umum/tidak spesifik (mis. cuma "
    "'analisis file ini' atau 'lihat file ini'), JANGAN cuma menjelaskan "
    "sekilas -- identifikasi jenis dokumennya lalu berikan analisis yang "
    "relevan untuk jenis dokumen itu (mis. rekening koran -> ekstrak "
    "transaksi jadi tabel + ringkasan saldo; data penjualan/invoice -> "
    "rekap total & rincian; bukti potong -> rekap jenis pajak & nilai "
    "potongan).\n\n"
    "Aturan wajib: JANGAN PERNAH mengarang angka, nama, atau data yang "
    "tidak benar-benar ada di file yang dikirim -- kalau ada bagian yang "
    "tidak terbaca jelas atau ambigu, sebutkan itu secara eksplisit alih-"
    "alih menebak. Jawab dalam Bahasa Indonesia formal, istilah akuntansi "
    "Indonesia yang baku (Debit/Kredit, bukan Debit/Credit), dan format "
    "tabel Markdown kalau user minta data ditabulasikan (mis. daftar "
    "transaksi)."
)

# [GABUNG] Konfigurasi retry (MAX_RETRY/RETRY_BASE_DELAY/RETRY_MAX_DELAY)
# SEKARANG ada di modules/claude_client.py (dipakai bersama semua
# pemanggil Claude, bisa dituning lewat env var CLAUDE_MAX_RETRY dst) --
# tidak didefinisikan ulang di sini, lihat _panggil_dengan_retry (alias
# import di atas).

# [BARU -- POIN 2: FILE BESAR/BINER] Batas ukuran file yang diparsing
# LOKAL (bukan limit Claude API -- ini limit kita sendiri, supaya server
# tidak coba muat file raksasa penuh ke memory). Cukup longgar utk file
# akuntansi wajar (rekap tahunan, dsb); kalau client benar-benar punya
# file lebih besar dari ini, sebaiknya dipecah dulu sebelum upload.
MAX_OFFICE_BYTES = 50 * 1024 * 1024    # 50 MB per xlsx/docx/pptx

# [BARU] Batas kasar utk validasi Content-Length SEBELUM body request
# dibaca penuh ke memory (lihat main.py: _cek_content_length_awal). Ambil
# limit TERBESAR di antara semua tipe file (Office 50MB) + buffer 10MB
# utk overhead header multipart/boundary -- bukan limit final per file
# (itu tetap divalidasi ulang lebih presisi oleh _validasi_pdf/_validasi_
# gambar/_validasi_office SETELAH tipe file diketahui dari nama file),
# murni pagar pertama biar upload yang JELAS kelewat besar ditolak tanpa
# perlu proses baca body request sama sekali.
MAX_UPLOAD_BYTES_PRACHECK = MAX_OFFICE_BYTES + 10 * 1024 * 1024
MAX_BARIS_SAMPEL_SHEET = 30            # baris SAMPEL yang diikutkan per sheet
                                        # (bukan seluruh isi -- lihat
                                        # _ekstrak_ringkasan_xlsx)
MAX_KARAKTER_CSV_UNTUK_RINGKASAN = 50_000  # CSV di atas ukuran ini dianggap
                                        # "dataset besar" -- diringkas pakai
                                        # pandas (shape/kolom/head/describe),
                                        # BUKAN dipotong mentah di karakter
                                        # ke-N seperti file teks biasa.

# [BARU -- POIN 2 TAMBAHAN] Sama seperti CSV di atas, tapi utk .json.
# Sebelumnya JSON besar TIDAK punya perlakuan khusus -- ikut jalur
# _potong_teks_jika_perlu (potong mentah di karakter ke-400.000), yang
# untuk JSON bisa berhenti di TENGAH struktur (mis. tengah objek/array),
# bikin AI kesulitan memahami bentuk datanya secara utuh.
MAX_KARAKTER_JSON_UNTUK_RINGKASAN = 50_000
MAX_KUNCI_JSON_DITAMPILKAN = 50  # utk JSON non-tabular (dict nested/kompleks),
                                        # jumlah maksimum kunci top-level yang
                                        # ditampilkan di ringkasan (lihat
                                        # _ringkas_dataset_json_besar)

# [BARU -- POIN 2, LANJUTAN] Cache hasil ekstraksi xlsx/xlsm/docx/xls/doc,
# analog persis dengan cache ekstraksi PDF di kertas_kerja.py (POIN 4:
# "gak proses ulang yang gak perlu"). Content-addressed pakai SHA-256 dari
# ISI FILE (bukan nama file) -- kalau 2 file beda nama tapi isinya identik
# byte-per-byte (mis. client upload ulang file yang sama), tetap cache-hit.
# Body hasil parsing (TANPA header "[Ringkasan file ... nama_file]") yang
# disimpan, supaya cache tetap valid dipakai lagi walau nama_file beda di
# panggilan berikutnya -- header ditempel belakangan pakai nama_file yang
# aktual saat itu.
_FOLDER_CACHE_EKSTRAKSI_OFFICE = Path(__file__).parent / "cache_ekstraksi_office"

# [BARU -- POIN 4, LANJUTAN] Cache untuk PDF, keyed by hash ISI file --
# beda dari cache office di atas (yang cache-nya berupa TEKS hasil
# ekstraksi), cache PDF di sini menyimpan base64 SIAP-KIRIM (encoding
# base64 itu sendiri bukan gratis -- untuk PDF besar, encode ulang tiap
# request itu kerja CPU yang terbuang percuma kalau file yang sama
# (misal client upload ulang rekening koran yang sama, atau retry dari
# frontend karena network drop) diproses lagi. TIDAK meng-cache jawaban
# Claude (isi PDF bisa ditanya hal beda-beda tiap kali), cuma meng-cache
# tahap encode/validasi yang deterministik terhadap isi file.
_FOLDER_CACHE_PDF = Path(__file__).parent / "cache_pdf_encoded"


def _hash_isi_file(isi_bytes: bytes) -> str:
    return hashlib.sha256(isi_bytes).hexdigest()


def _path_cache_pdf(hash_isi: str) -> Path:
    _FOLDER_CACHE_PDF.mkdir(exist_ok=True)
    return _FOLDER_CACHE_PDF / f"{hash_isi}.b64"


def _muat_cache_pdf_base64(hash_isi: str) -> Optional[str]:
    """Baca base64 PDF dari cache disk (bukan pickle -- teks base64 murni,
    lebih ringan & bisa dibaca langsung). None kalau cache-miss/rusak,
    cache tidak pernah melempar error (pola sama dgn cache office)."""
    path = _path_cache_pdf(hash_isi)
    if not path.exists():
        return None
    try:
        data = path.read_text(encoding="ascii")
        try:
            os.utime(path, None)  # tandai "baru diakses" utk cache_cleanup.py
        except OSError:
            pass
        return data
    except Exception as e:  # noqa: BLE001
        logger.warning(f"⚠️ Gagal baca cache PDF ({path.name}): {e} -- dianggap cache-miss.")
        return None


def _simpan_cache_pdf_base64(hash_isi: str, base64_data: str) -> None:
    path = _path_cache_pdf(hash_isi)
    try:
        tmp = path.with_suffix(".tmp")
        tmp.write_text(base64_data, encoding="ascii")
        tmp.replace(path)
    except Exception as e:  # noqa: BLE001
        logger.warning(f"⚠️ Gagal simpan cache PDF ke {path.name}: {e} -- dilewati, tidak fatal.")


def _encode_pdf_base64_dengan_cache(isi_bytes: bytes) -> str:
    """[BARU -- POIN 4] Encode PDF ke base64, cache-hit kalau isi file
    (byte-per-byte) sudah pernah di-encode sebelumnya -- menghindari
    kerja CPU encode ulang utk PDF besar yang sama (upload ulang, retry
    frontend, dsb). Validasi ukuran/halaman TETAP dijalankan tiap kali
    di caller (_validasi_pdf) -- cache HANYA utk tahap encode, bukan
    utk melewati validasi."""
    hash_isi = _hash_isi_file(isi_bytes)
    cached = _muat_cache_pdf_base64(hash_isi)
    if cached is not None:
        return cached
    encoded = base64.b64encode(isi_bytes).decode("utf-8")
    _simpan_cache_pdf_base64(hash_isi, encoded)
    return encoded


def _path_cache_ekstraksi_office(hash_isi: str) -> Path:
    _FOLDER_CACHE_EKSTRAKSI_OFFICE.mkdir(exist_ok=True)
    return _FOLDER_CACHE_EKSTRAKSI_OFFICE / f"{hash_isi}.pkl"


def _muat_cache_ekstraksi_office(hash_isi: str) -> Optional[str]:
    """Baca body hasil ekstraksi dari cache disk. Return None kalau
    cache-miss ATAU cache rusak/tidak terbaca -- cache murni optimasi,
    TIDAK PERNAH melempar error ke pemanggil (sama pola dgn kertas_kerja.py)."""
    path = _path_cache_ekstraksi_office(hash_isi)
    if not path.exists():
        return None
    try:
        with open(path, "rb") as f:
            data = pickle.load(f)
        body = data.get("body")
        # [BARU -- TTL/LRU] Sentuh mtime saat cache-hit -- lihat catatan
        # identik di kertas_kerja.py::_muat_cache_ekstraksi_pdf. Dipakai
        # modules/cache_cleanup.py sebagai proxy "terakhir diakses".
        try:
            os.utime(path, None)
        except OSError:
            pass
        return body
    except Exception as e:  # noqa: BLE001
        logger.warning(f"⚠️ Gagal baca cache ekstraksi office ({path.name}): {e} -- dianggap cache-miss.")
        return None


def _simpan_cache_ekstraksi_office(hash_isi: str, body: str) -> None:
    """Simpan body hasil ekstraksi ke cache disk. Tulis ke .tmp dulu lalu
    rename (atomic) -- menghindari proses lain baca file cache yang
    setengah tertulis kalau proses ini terhenti di tengah jalan. Gagal
    simpan cache TIDAK BOLEH menjatuhkan proses ekstraksi utama."""
    path = _path_cache_ekstraksi_office(hash_isi)
    try:
        tmp = path.with_suffix(".tmp")
        with open(tmp, "wb") as f:
            pickle.dump({"body": body}, f)
        tmp.replace(path)
    except Exception as e:  # noqa: BLE001
        logger.warning(f"⚠️ Gagal menyimpan cache ekstraksi office ke {path.name}: {e} -- dilewati, tidak fatal.")


def _hitung_perkiraan_halaman_pdf(isi_pdf: bytes) -> Optional[int]:
    """
    [BARU] Hitung PERKIRAAN jumlah halaman PDF tanpa dependency baru
    (tidak pakai pypdf/PyPDF2) -- cukup untuk validasi limit 100 halaman
    di sini, BUKAN untuk kebutuhan presisi (mis. penomoran halaman).

    Menghitung kemunculan pola objek "/Type /Page" (bukan "/Pages", yang
    beda objek -- lihat spec PDF) di raw bytes file. Ini heuristik yang
    berhasil untuk PDF standar (non-terenkripsi, non-linierized eksotis)
    -- untuk kasus PDF yang strukturnya tidak umum, return None (caller
    memperlakukan ini sebagai "tidak bisa dipastikan", BUKAN gagal validasi
    -- lebih baik lolos & biarkan API sendiri yang menolak kalau ternyata
    kepanjangan, daripada salah menolak PDF yang sah).
    """
    try:
        cocok = re.findall(rb"/Type\s*/Page[^s]", isi_pdf)
        jumlah = len(cocok)
        return jumlah if jumlah > 0 else None
    except Exception:  # noqa: BLE001 -- heuristik, tidak boleh melempar error
        return None


def _validasi_pdf(isi_bytes: bytes, nama_file: str) -> None:
    """[BARU] Validasi PDF sebelum dikirim ke API -- melempar ValueError
    dengan pesan jelas kalau melanggar limit resmi, supaya caller (endpoint)
    bisa balikin 400 yang informatif, bukan menunggu API menolak dengan
    pesan generik setelah upload base64 selesai (buang-buang bandwidth &
    waktu untuk file yang sudah pasti gagal)."""
    ukuran = len(isi_bytes)
    if ukuran > MAX_PDF_BYTES:
        raise ValueError(
            f"'{nama_file}' berukuran {ukuran / (1024*1024):.1f} MB, melebihi "
            f"limit Claude API {MAX_PDF_BYTES // (1024*1024)} MB per PDF. "
            "Pecah file jadi beberapa bagian lebih kecil."
        )
    jumlah_halaman = _hitung_perkiraan_halaman_pdf(isi_bytes)
    if jumlah_halaman is not None and jumlah_halaman > MAX_PDF_HALAMAN:
        raise ValueError(
            f"'{nama_file}' berisi kira-kira {jumlah_halaman} halaman, melebihi "
            f"limit Claude API {MAX_PDF_HALAMAN} halaman per request untuk analisis "
            "visual penuh. Pecah file jadi beberapa bagian per rentang halaman."
        )


def _validasi_gambar(isi_bytes: bytes, nama_file: str) -> None:
    """[BARU] Validasi gambar sebelum dikirim -- limit resmi 5MB per gambar."""
    ukuran = len(isi_bytes)
    if ukuran > MAX_GAMBAR_BYTES:
        raise ValueError(
            f"'{nama_file}' berukuran {ukuran / (1024*1024):.1f} MB, melebihi "
            f"limit Claude API {MAX_GAMBAR_BYTES // (1024*1024)} MB per gambar. "
            "Kompres atau perkecil resolusi gambar dulu."
        )


# ============================================================
# [BARU -- SEMENTARA] EKSTRAKSI TEKS PDF (bukan dikirim base64 mentah)
#
# Alasan: kirim PDF sebagai content block "document" (lihat kirim_pdf_ke_ai
# & siapkan_konten_pesan_dari_file versi lama) HANYA bisa dibaca oleh
# Claude (Anthropic Messages API) -- Groq/model OpenAI-compatible lain
# TIDAK punya fitur "document understanding" tsb, cuma terima teks polos
# di chat.completions.create(). Selama ANTHROPIC_API_KEY belum ada
# saldo/belum aktif, PDF di sini SEMENTARA diekstrak jadi teks UTUH dulu
# (pakai pdfplumber, library yg sama yg sudah dipakai akuntansi_ai.py
# utk baca rekening koran -- lihat _baca_pdf_sebagai_lembar di sana),
# baru teks hasil ekstraksi itu yang dikirim ke provider chat (Claude
# kalau ada key, fallback Groq kalau tidak -- lihat
# _konfigurasi_provider_file_reader di bawah).
#
# CATATAN: ini mengorbankan kemampuan baca LAYOUT/VISUAL langsung
# (tabel rumit, PDF hasil scan/gambar tanpa teks asli tidak akan
# terbaca sama sekali oleh ekstraksi teks polos ini) -- begitu
# ANTHROPIC_API_KEY sudah aktif lagi, tinggal balikin PDF ke jalur
# base64 "document" seperti semula (lihat siapkan_konten_pesan_dari_file).
# ============================================================

def _ekstrak_teks_pdf_penuh(isi_bytes: bytes, nama_file: str) -> str:
    """Ekstrak SELURUH teks PDF (semua halaman, digabung apa adanya per
    halaman) pakai pdfplumber -- TIDAK diringkas/dipotong di sini (lihat
    _potong_teks_jika_perlu yang dipanggil terpisah oleh caller kalau
    hasilnya kepanjangan). Melempar ValueError kalau PDF gagal dibuka
    (mis. terenkripsi/password, korup) atau ternyata tidak ada teks sama
    sekali (mis. PDF hasil scan murni tanpa OCR -- di luar cakupan
    ekstraksi teks polos ini)."""
    try:
        import pdfplumber
    except ImportError as e:
        raise RuntimeError(
            "Gagal membaca file PDF -- library 'pdfplumber' belum ter-install. "
            "Jalankan: pip install pdfplumber --break-system-packages"
        ) from e

    bagian_per_halaman: List[str] = []
    try:
        with pdfplumber.open(io.BytesIO(isi_bytes)) as pdf:
            for i_hal, page in enumerate(pdf.pages, start=1):
                teks_halaman = page.extract_text() or ""
                bagian_per_halaman.append(f"--- Halaman {i_hal} ---\n{teks_halaman}")
    except Exception as e:  # noqa: BLE001 -- mis. PDF terenkripsi/korup
        raise ValueError(
            f"Gagal mengekstrak teks dari '{nama_file}': {e}. "
            "Kemungkinan file terenkripsi/berpassword atau formatnya tidak standar."
        ) from e

    teks_gabungan = "\n\n".join(bagian_per_halaman).strip()
    if not teks_gabungan:
        raise ValueError(
            f"'{nama_file}' tidak mengandung teks yang bisa diekstrak (kemungkinan "
            "PDF hasil scan/gambar tanpa OCR) -- jalur ekstraksi teks sementara ini "
            "tidak bisa membacanya. Butuh ANTHROPIC_API_KEY aktif utk baca PDF jenis ini."
        )
    return teks_gabungan


# [BARU -- SEMENTARA] Provider chat KHUSUS utk semua tipe file yang
# CONTENT-nya sudah berupa TEKS POLOS (teks/csv/json/html, hasil ekstraksi
# PDF, hasil ekstraksi Office xlsx/xlsm/docx/pptx/xls/doc) -- SAMA PERSIS
# pola prioritasnya dengan akuntansi_ai._konfigurasi_provider_kategorisasi()
# (Claude dulu kalau ANTHROPIC_API_KEY terisi, fallback Groq kategorisasi/
# GROQ_API_KEY_KATEGORISASI kalau tidak) -- SENGAJA duplikasi kecil di sini
# (bukan import dari akuntansi_ai.py) supaya modul ini tetap bisa dipakai
# berdiri sendiri tanpa import silang antar modul yang tidak perlu.
#
# TIDAK dipakai untuk tipe "gambar" -- gambar butuh model VISION (bisa
# "melihat" isi gambar), dan model Groq yang dikonfigurasi di sini
# ("openai/gpt-oss-120b" / GROQ_MODEL_KATEGORISASI) adalah model TEKS
# MURNI, tidak punya kemampuan vision. Gambar tetap WAJIB lewat Claude
# (lihat kirim_file_ke_ai/kirim_file_ke_ai_stream -- dibedakan berdasar
# tipe content: string = teks polos -> boleh fallback, list = ada content
# block gambar -> Claude saja).
def _konfigurasi_provider_file_reader() -> list[dict]:
    """Balikin LIST provider utk membalas pertanyaan atas content TEKS
    (bukan gambar), urut prioritas: Claude dulu, Groq (kategorisasi) sbg
    fallback. List kosong kalau tidak ada key sama sekali."""
    daftar = []
    claude_key = os.environ.get("ANTHROPIC_API_KEY")
    if claude_key:
        daftar.append({"tipe": "anthropic", "api_key": claude_key, "model": MODEL_DEFAULT, "nama": "Claude"})
    groq_key = os.environ.get("GROQ_API_KEY_KATEGORISASI") or os.environ.get("GROQ_API_KEY")
    if groq_key:
        daftar.append({
            "tipe": "openai_compatible",
            "api_key": groq_key,
            "base_url": "https://api.groq.com/openai/v1",
            "model": os.environ.get("GROQ_MODEL_KATEGORISASI", "openai/gpt-oss-120b"),
            "nama": "Groq (fallback baca PDF)",
            "extra_params": {"reasoning_effort": "medium"},
        })
    return daftar


def _tanya_teks_ke_provider(prompt: str) -> str:
    """[BARU -- SEMENTARA] Kirim `prompt` teks polos (isi file teks/PDF/
    Office hasil ekstraksi + pertanyaan) ke provider pertama yang key-nya
    terisi (lihat _konfigurasi_provider_file_reader), coba provider
    berikutnya kalau satu gagal. Melempar RuntimeError kalau tidak ada
    provider sama sekali, atau exception terakhir kalau SEMUA provider
    di list gagal."""
    daftar_provider = _konfigurasi_provider_file_reader()
    if not daftar_provider:
        raise RuntimeError(
            "Tidak ada API key aktif untuk membaca file (ANTHROPIC_API_KEY / "
            "GROQ_API_KEY_KATEGORISASI / GROQ_API_KEY semuanya kosong)."
        )

    error_terakhir: Optional[Exception] = None
    for konfig in daftar_provider:
        try:
            if konfig["tipe"] == "anthropic":
                response = _panggil_dengan_retry(
                    model=konfig["model"],
                    max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
                    system=SYSTEM_PROMPT_FILE_READER,
                    messages=[{"role": "user", "content": prompt}],
                )
                return response.content[0].text
            else:  # "openai_compatible" -- Groq
                # [BARU] API OpenAI-compatible (Groq) TIDAK punya parameter
                # `system` terpisah seperti Anthropic -- system prompt
                # dikirim sbg pesan pertama ber-role "system" di list
                # `messages` yang sama.
                import openai
                client = openai.OpenAI(api_key=konfig["api_key"], base_url=konfig["base_url"], max_retries=0)
                response = client.chat.completions.create(
                    model=konfig["model"],
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT_FILE_READER},
                        {"role": "user", "content": prompt},
                    ],
                    **konfig.get("extra_params", {}),
                )
                return response.choices[0].message.content
        except Exception as e:  # noqa: BLE001 -- coba provider berikutnya di list
            logger.warning(f"⚠️ Provider '{konfig['nama']}' gagal baca file: {e} -- coba provider berikutnya.")
            error_terakhir = e
            continue

    raise error_terakhir or RuntimeError("Gagal memanggil semua provider file reader.")


# ============================================================
# A. FILE TEKS (.md, .txt, .csv) -- dikirim sebagai teks biasa
# ============================================================

def kirim_file_teks_ke_ai(path_file: str | Path, pertanyaan: str, model: str = MODEL_DEFAULT) -> str:
    """Baca file teks apa adanya, tempel ke prompt -- tidak perlu parsing khusus."""
    isi_file = Path(path_file).read_text(encoding="utf-8")
    isi_file, dipotong = _potong_teks_jika_perlu(isi_file)

    prompt = f"Berikut isi file:\n\n{isi_file}\n\n{pertanyaan}"
    if dipotong:
        prompt += "\n\n(Catatan: isi file dipotong karena terlalu panjang -- lihat bagian awal saja.)"

    response = _panggil_dengan_retry(
        model=model,
        max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
        system=SYSTEM_PROMPT_FILE_READER,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.content[0].text


# ============================================================
# B. GAMBAR (.png, .jpg) -- dikirim sebagai base64 content block
# ============================================================

_MEDIA_TYPE_GAMBAR = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


def kirim_gambar_ke_ai(path_gambar: str | Path, pertanyaan: str, model: str = MODEL_DEFAULT) -> str:
    """Kirim file gambar sebagai base64 -- Claude langsung "melihat" isinya
    (termasuk teks di dalam gambar, tabel, layout, dll)."""
    path_gambar = Path(path_gambar)
    media_type = _MEDIA_TYPE_GAMBAR.get(path_gambar.suffix.lower())
    if not media_type:
        raise ValueError(f"Ekstensi gambar tidak didukung: {path_gambar.suffix}")

    isi_bytes = path_gambar.read_bytes()
    _validasi_gambar(isi_bytes, path_gambar.name)
    base64_data = base64.b64encode(isi_bytes).decode("utf-8")

    response = _panggil_dengan_retry(
        model=model,
        max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
        system=SYSTEM_PROMPT_FILE_READER,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {"type": "base64", "media_type": media_type, "data": base64_data},
                },
                {"type": "text", "text": pertanyaan},
            ],
        }],
    )
    return response.content[0].text


# ============================================================
# C. PDF -- dikirim sebagai base64 content block tipe "document"
# ============================================================

def kirim_pdf_ke_ai(path_pdf: str | Path, pertanyaan: str, model: str = MODEL_DEFAULT) -> str:
    """Kirim file PDF sebagai base64 -- Claude langsung membaca teks &
    layout-nya, TANPA perlu pdfplumber/pytesseract di sisi kita."""
    path_pdf = Path(path_pdf)
    isi_bytes = path_pdf.read_bytes()
    _validasi_pdf(isi_bytes, path_pdf.name)
    base64_data = _encode_pdf_base64_dengan_cache(isi_bytes)

    response = _panggil_dengan_retry(
        model=model,
        max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
        system=SYSTEM_PROMPT_FILE_READER,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "document",
                    "source": {"type": "base64", "media_type": "application/pdf", "data": base64_data},
                },
                {"type": "text", "text": pertanyaan},
            ],
        }],
    )
    return response.content[0].text


# ============================================================
# D. Versi terima BYTES langsung + auto-deteksi tipe file --
# dipakai dari FastAPI (UploadFile.read() menghasilkan bytes di memori,
# BUKAN path di disk -- sama pola seperti daftar_file_pdf di kertas_kerja.py)
# ============================================================

_EKSTENSI_TEKS = {".md", ".txt", ".csv", ".json", ".html"}
_EKSTENSI_GAMBAR = set(_MEDIA_TYPE_GAMBAR.keys())
# [BARU -- POIN 2] .xlsx/.xlsm/.docx TIDAK bisa dikirim mentah ke Claude
# API (bukan format yang didukung content block apa pun di API) -- jadi
# harus diparsing lokal dulu (lihat _ekstrak_ringkasan_xlsx/_ekstrak_teks_docx)
# sebelum hasil ekstraksinya (berupa teks) dikirim ke Claude.
_EKSTENSI_OFFICE = {".xlsx", ".xlsm", ".docx", ".pptx"}
# [BARU -- POIN 2, LANJUTAN] Format Office LAMA (binary/OLE, pra-2007) --
# .xls TIDAK bisa dibuka openpyxl (openpyxl cuma paham format .xlsx berbasis
# XML/zip), .doc TIDAK bisa dibuka python-docx (sama alasan). Perlu library
# beda (xlrd utk .xls, textract/antiword utk .doc) -- dipisah jadi kategori
# sendiri supaya pesan error saat library-nya belum terpasang jelas
# menyebut file LAMA ini, bukan tercampur dgn error xlsx/docx modern.
_EKSTENSI_OFFICE_LAMA = {".xls", ".doc"}


def _ekstrak_body_pptx(isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU] Parsing LOKAL file .pptx pakai python-pptx -- BUKAN dikirim
    mentah ke Claude API (bukan format yang didukung content block apa
    pun, sama seperti xlsx/docx). Diambil per slide: judul, semua text
    frame (bullet points dst), isi tabel, DAN speaker notes (sering
    berisi konteks penting yang tidak kelihatan di slide itu sendiri).

    TIDAK mengekstrak gambar/chart embedded (analog keterbatasan yang
    sama di _ekstrak_body_docx/_ekstrak_body_xlsx) -- kalau slide isinya
    murni gambar/screenshot tanpa teks, bagian itu akan kosong di hasil
    ekstraksi.

    Import python-pptx LAZY -- sama pola dgn python-docx/xlrd/textract.
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "Library 'python-pptx' belum terinstall -- jalankan "
            "`pip install python-pptx` untuk mendukung upload file .pptx."
        ) from e

    try:
        prs = Presentation(io.BytesIO(isi_bytes))
    except Exception as e:  # noqa: BLE001
        raise ValueError(f"Gagal membuka '{nama_file}' sebagai file PowerPoint: {e}")

    bagian: List[str] = [f"{len(prs.slides)} slide"]
    for i, slide in enumerate(prs.slides, start=1):
        bagian.append(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_table:
                tabel = shape.table
                for row in tabel.rows:
                    bagian.append(" | ".join(cell.text.strip() for cell in row.cells))
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    teks_para = "".join(run.text for run in para.runs).strip()
                    if teks_para:
                        bagian.append(teks_para)

        if slide.has_notes_slide:
            teks_catatan = slide.notes_slide.notes_text_frame.text.strip()
            if teks_catatan:
                bagian.append(f"[Catatan pembicara] {teks_catatan}")

    return "\n".join(bagian)


def _ekstrak_teks_pptx(isi_bytes: bytes, nama_file: str) -> Tuple[str, bool]:
    """[BARU] Wrapper ber-cache di atas _ekstrak_body_pptx, pola identik
    dengan _ekstrak_ringkasan_xlsx/_ekstrak_teks_docx. Return
    (teks_lengkap, dari_cache).

    [FIX -- BUG timeout/recovery hilang di jalur satu-file] Sama seperti
    _ekstrak_ringkasan_xlsx -- lihat catatan lengkap di sana."""
    hash_isi = _hash_isi_file(isi_bytes)
    body = _muat_cache_ekstraksi_office(hash_isi)
    dari_cache = body is not None
    if body is None:
        body = _jalankan_ekstraksi_di_proses(".pptx", isi_bytes, nama_file)
        _simpan_cache_ekstraksi_office(hash_isi, body)

    teks_lengkap = f"[Isi presentasi PowerPoint '{nama_file}'{' -- dari cache' if dari_cache else ''}] {body}"
    return teks_lengkap, dari_cache


def deteksi_tipe_file(nama_file: str) -> str:
    """Return 'teks' | 'gambar' | 'pdf' | 'office' | 'office_lama' |
    'tidak_didukung' dari ekstensi nama file. 'office' = xlsx/xlsm/docx
    (format XML/zip modern, lihat _EKSTENSI_OFFICE); 'office_lama' =
    xls/doc (format biner lama, lihat _EKSTENSI_OFFICE_LAMA) -- keduanya
    perlu diparsing lokal, tidak pernah dikirim mentah ke Claude API."""
    ext = Path(nama_file).suffix.lower()
    if ext in _EKSTENSI_TEKS:
        return "teks"
    if ext in _EKSTENSI_GAMBAR:
        return "gambar"
    if ext == ".pdf":
        return "pdf"
    if ext in _EKSTENSI_OFFICE:
        return "office"
    if ext in _EKSTENSI_OFFICE_LAMA:
        return "office_lama"
    return "tidak_didukung"


def _ekstrak_body_xls_lama(isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU -- POIN 2, LANJUTAN] Parsing LOKAL file .xls (format biner lama,
    BIFF -- BUKAN xlsx/XML) pakai xlrd. openpyxl TIDAK bisa membuka format
    ini sama sekali (beda struktur file total dari xlsx), makanya perlu
    library terpisah. xlrd versi modern (>=2.0) SENGAJA sudah membuang
    dukungan xlsx (murni fokus .xls lama) -- jangan salah pasang xlrd utk
    file .xlsx, itu tugas openpyxl (lihat _ekstrak_body_xlsx).

    Import xlrd LAZY (sama pola dgn python-docx di _ekstrak_body_docx) --
    modul ini tetap bisa dipakai tanpa xlrd terpasang kalau client tidak
    pernah upload .xls lama.
    """
    try:
        import xlrd
    except ImportError as e:
        raise RuntimeError(
            "Library 'xlrd' belum terinstall -- jalankan `pip install xlrd` "
            "untuk mendukung upload file .xls (format Excel lama). Alternatif "
            "lebih cepat: minta client export ulang sebagai .xlsx."
        ) from e

    try:
        wb = xlrd.open_workbook(file_contents=isi_bytes)
    except Exception as e:  # noqa: BLE001
        raise ValueError(f"Gagal membuka '{nama_file}' sebagai file Excel lama (.xls): {e}")

    bagian: List[str] = [f"{wb.nsheets} sheet (format .xls lama)"]
    for ws in wb.sheets():
        total_baris = ws.nrows
        n_sampel = min(total_baris, MAX_BARIS_SAMPEL_SHEET)
        bagian.append(f"\n=== Sheet '{ws.name}' ({total_baris} baris total) ===")
        for i in range(n_sampel):
            nilai_baris = ws.row_values(i)
            teks_baris = " | ".join("" if v is None else str(v) for v in nilai_baris)
            bagian.append(teks_baris)
        if total_baris > MAX_BARIS_SAMPEL_SHEET:
            bagian.append(f"... ({total_baris - MAX_BARIS_SAMPEL_SHEET} baris lain tidak ditampilkan, hanya sampel)")

    return "\n".join(bagian)


def _ekstrak_body_doc_lama(isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU -- POIN 2, LANJUTAN] Parsing LOKAL file .doc (format biner lama,
    OLE Compound File -- BUKAN docx/XML) pakai textract. python-docx TIDAK
    bisa membuka format ini sama sekali.

    Import textract LAZY -- library ini SENGAJA tidak dijadikan dependency
    wajib (berat, & di banyak OS masih butuh binary eksternal 'antiword'
    terpasang di sistem juga) supaya modul ini tetap ringan dipasang kalau
    client tidak pernah kirim .doc lama (kasus jarang dibanding .docx).
    """
    try:
        import textract
    except ImportError as e:
        raise RuntimeError(
            "Library 'textract' (+ binary 'antiword' di sistem) belum "
            "terinstall -- jalankan `pip install textract` untuk mendukung "
            "upload file .doc (format Word lama). Alternatif lebih cepat: "
            "minta client export ulang sebagai .docx."
        ) from e

    try:
        with tempfile_office_lama(isi_bytes, ".doc") as path_temp:
            teks = textract.process(str(path_temp)).decode("utf-8", errors="replace")
    except Exception as e:  # noqa: BLE001
        raise ValueError(f"Gagal membuka '{nama_file}' sebagai file Word lama (.doc): {e}")

    return teks


def tempfile_office_lama(isi_bytes: bytes, suffix: str):
    """[BARU] textract butuh PATH FILE DI DISK (tidak terima bytes langsung
    seperti openpyxl/python-docx) -- helper kecil ini bungkus jadi context
    manager file temp yang otomatis terhapus, supaya tidak ada file
    /tmp/*.doc client yang tertinggal setelah selesai."""
    import contextlib
    import tempfile

    @contextlib.contextmanager
    def _ctx():
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(isi_bytes)
            path_temp = Path(f.name)
        try:
            yield path_temp
        finally:
            path_temp.unlink(missing_ok=True)

    return _ctx()


def ekstrak_office_lama(isi_bytes: bytes, nama_file: str) -> Tuple[str, bool]:
    """
    [BARU -- POIN 2, LANJUTAN] Titik masuk untuk tipe 'office_lama' (.xls/.doc)
    -- dispatch ke _ekstrak_body_xls_lama/_ekstrak_body_doc_lama sesuai
    ekstensi, dengan cache disk yang sama persis polanya dengan
    _ekstrak_ringkasan_xlsx/_ekstrak_teks_docx (POIN 4). Return
    (teks_lengkap, dari_cache).

    [FIX -- BUG timeout/recovery hilang di jalur satu-file] Sama seperti
    _ekstrak_ringkasan_xlsx -- delegasikan ke _jalankan_ekstraksi_di_proses
    (proses terpisah + timeout + recovery), bukan panggil _ekstrak_body_*
    langsung. .xls lewat xlrd murni Python (GIL-bound, rawan hang di
    struktur BIFF rusak) & .doc lewat textract (subprocess antiword) --
    keduanya sama-sama butuh proteksi ini, bukan cuma .xlsx/.docx.
    """
    ext = Path(nama_file).suffix.lower()
    if ext not in (".xls", ".doc"):
        raise ValueError(f"Ekstensi office lama tidak dikenali: '{ext}'")
    hash_isi = _hash_isi_file(isi_bytes)
    body = _muat_cache_ekstraksi_office(hash_isi)
    dari_cache = body is not None
    if body is None:
        body = _jalankan_ekstraksi_di_proses(ext, isi_bytes, nama_file)
        _simpan_cache_ekstraksi_office(hash_isi, body)

    label = "Excel" if ext == ".xls" else "Word"
    teks_lengkap = f"[Isi dokumen {label} lama '{nama_file}'{' -- dari cache' if dari_cache else ''}]\n{body}"
    return teks_lengkap, dari_cache


def _validasi_office(isi_bytes: bytes, nama_file: str) -> None:
    """[BARU] Validasi ukuran xlsx/docx sebelum diparsing lokal -- limit
    kita sendiri (bukan limit Claude API, karena file ini tidak pernah
    dikirim mentah ke API), supaya server tidak coba muat file raksasa
    penuh ke memory sekaligus."""
    ukuran = len(isi_bytes)
    if ukuran > MAX_OFFICE_BYTES:
        raise ValueError(
            f"'{nama_file}' berukuran {ukuran / (1024*1024):.1f} MB, melebihi "
            f"limit {MAX_OFFICE_BYTES // (1024*1024)} MB untuk file xlsx/docx yang "
            "diparsing di server ini. Pecah file jadi beberapa bagian lebih kecil."
        )


# [BARU -- POIN 3] Signature byte OLE Compound File (dulu dipakai format
# Excel/Word BINER lama .xls/.doc, tapi juga dipakai sbg "kontainer" utk
# xlsx/docx MODERN yang dienkripsi password -- spec MS-OFFCRYPTO: file
# xlsx/docx yang dikunci password TIDAK disimpan sbg zip/OOXML biasa
# (yang mulai dgn "PK\x03\x04"), melainkan dibungkus dalam OLE Compound
# File berisi stream "EncryptedPackage". Jadi kalau ekstensi file .xlsx/
# .xlsm/.docx TAPI byte awalnya cocok signature ini (bukan "PK"), file
# itu HAMPIR PASTI terkunci password -- dideteksi SEBELUM openpyxl/
# python-docx dipanggil, supaya pesan errornya jelas ("terkunci
# password") alih-alih exception generik (BadZipFile dsb) yang
# membingungkan user.
_SIGNATURE_OLE_COMPOUND_FILE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _tampak_seperti_file_terkunci_password(isi_bytes: bytes) -> bool:
    """[BARU -- POIN 3] True kalau byte awal file cocok signature OLE
    Compound File -- dipakai HANYA utk ekstensi xlsx/xlsm/docx (format
    modern yang seharusnya zip/"PK", bukan OLE) sebagai sinyal kuat file
    itu dienkripsi password, sebelum openpyxl/python-docx dipanggil."""
    return isi_bytes[:8] == _SIGNATURE_OLE_COMPOUND_FILE


def _pesan_error_kemungkinan_password(nama_file: str, label_format: str, error_asli: Exception) -> ValueError:
    """[BARU -- POIN 3] Bungkus exception generik (BadZipFile dkk) jadi
    ValueError dengan pesan jelas & actionable, dipakai sbg fallback di
    except block _ekstrak_body_xlsx/_ekstrak_body_docx utk kasus yang
    TIDAK ketangkap deteksi signature di atas (mis. file korup dgn cara
    lain yang kebetulan gejalanya mirip) -- tetap menyebut kemungkinan
    password sbg salah satu penyebab paling umum, tanpa mengklaim pasti."""
    return ValueError(
        f"'{nama_file}' gagal dibuka sebagai file {label_format} -- kemungkinan besar file ini "
        f"TERKUNCI PASSWORD (proteksi buka-file, bukan sekadar read-only). Buka file ini secara "
        "manual, masukkan passwordnya, lalu simpan ulang TANPA proteksi password (Save As / "
        f"Remove Password) sebelum upload lagi. (Detail teknis: {type(error_asli).__name__}: {error_asli})"
    )


def _ekstrak_body_xlsx(isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU -- POIN 2] Parsing LOKAL file .xlsx/.xlsm pakai openpyxl (analog
    "sandbox Linux + pandas/openpyxl" yang dipakai Claude sendiri utk file
    besar/biner) -- BUKAN dikirim mentah ke Claude API (tidak didukung).

    [DIPERBAIKI] Sebelumnya `list(ws.iter_rows(values_only=True))` memuat
    SELURUH baris ke memory dulu baru dipotong sampel -- ini MENIADAKAN
    manfaat `read_only=True` (streaming) di atas untuk sheet besar (puluhan
    ribu baris data akuntansi setahun bisa berat di memory). Sekarang
    di-iterate SEKALI secara streaming: baris disimpan ke `sampel` HANYA
    selama belum mencapai MAX_BARIS_SAMPEL_SHEET, baris setelahnya cuma
    dihitung (`total_baris += 1`) lalu DIBUANG, tidak pernah disimpan --
    jejak memory tetap O(MAX_BARIS_SAMPEL_SHEET), bukan O(total baris),
    walau tetap perlu 1x lintasan penuh untuk tahu total baris sungguhan.

    Return BODY saja (tanpa header "[Ringkasan file ...]") -- header
    ditempel di pemanggil (lihat _ekstrak_ringkasan_xlsx) supaya body ini
    valid dicache lepas dari nama_file (lihat _FOLDER_CACHE_EKSTRAKSI_OFFICE).
    """
    # [BARU -- POIN 3] Cek signature SEBELUM openpyxl dipanggil -- kalau
    # cocok OLE Compound File, langsung pesan jelas "terkunci password"
    # tanpa perlu tunggu openpyxl melempar BadZipFile yang membingungkan.
    if _tampak_seperti_file_terkunci_password(isi_bytes):
        raise ValueError(
            f"'{nama_file}' terkunci password (proteksi buka-file) -- tidak bisa dibaca otomatis. "
            "Buka file ini secara manual, masukkan passwordnya, lalu simpan ulang (Save As) TANPA "
            "proteksi password sebelum upload lagi."
        )
    try:
        wb = openpyxl.load_workbook(io.BytesIO(isi_bytes), read_only=True, data_only=True)
    except Exception as e:  # noqa: BLE001
        raise _pesan_error_kemungkinan_password(nama_file, "Excel", e)

    bagian: List[str] = [f"{len(wb.sheetnames)} sheet"]
    for nama_sheet in wb.sheetnames:
        ws = wb[nama_sheet]
        sampel: List[tuple] = []
        total_baris = 0
        for baris in ws.iter_rows(values_only=True):
            total_baris += 1
            if len(sampel) < MAX_BARIS_SAMPEL_SHEET:
                sampel.append(baris)
            # baris di luar sampel SENGAJA tidak disimpan kemana pun --
            # cuma numpang lewat buat hitung total_baris, langsung dibuang
            # (di-garbage-collect) begitu iterasi lanjut ke baris berikutnya.

        bagian.append(f"\n=== Sheet '{nama_sheet}' ({total_baris} baris total) ===")
        for baris in sampel:
            teks_baris = " | ".join("" if v is None else str(v) for v in baris)
            bagian.append(teks_baris)
        if total_baris > MAX_BARIS_SAMPEL_SHEET:
            bagian.append(f"... ({total_baris - MAX_BARIS_SAMPEL_SHEET} baris lain tidak ditampilkan, hanya sampel)")

    wb.close()
    return "\n".join(bagian)


def _ekstrak_ringkasan_xlsx(isi_bytes: bytes, nama_file: str) -> Tuple[str, bool]:
    """
    [BARU -- POIN 2, LANJUTAN] Wrapper ber-cache di atas _ekstrak_body_xlsx --
    cek cache disk dulu (keyed by hash isi file) sebelum parsing openpyxl
    yang sebenarnya (POIN 4: "gak proses ulang yang gak perlu", sama pola
    dgn cache ekstraksi PDF di kertas_kerja.py). Return (teks_lengkap, dari_cache).

    [FIX -- BUG: TIMEOUT/RECOVERY HILANG DI JALUR SATU-FILE] Sebelumnya
    cache-miss di sini memanggil `_ekstrak_body_xlsx()` LANGSUNG di thread
    pemanggil -- BEDA dari jalur multi-file (_olah_satu_file_untuk_multi)
    yang sudah lebih dulu didelegasikan ke `_jalankan_ekstraksi_di_proses()`
    (proses terpisah + timeout 120 detik + recovery BrokenProcessPool).
    Akibatnya file .xlsx korup/struktur aneh yang membuat openpyxl hang
    di endpoint SATU FILE (/api/ai-baca-file -- jalur PALING SERING dipakai)
    bisa menggantung SELAMANYA tanpa proteksi apa pun, dan karena endpoint
    itu jalan lewat asyncio.to_thread() (thread pool default asyncio yang
    ukurannya terbatas), beberapa upload macet beruntun bisa menghabiskan
    seluruh thread pool & ikut memblokir endpoint lain. Sekarang disamakan:
    LEWAT proses terpisah juga, persis pola jalur multi-file.
    """
    hash_isi = _hash_isi_file(isi_bytes)
    body = _muat_cache_ekstraksi_office(hash_isi)
    dari_cache = body is not None
    if body is None:
        body = _jalankan_ekstraksi_di_proses(Path(nama_file).suffix.lower(), isi_bytes, nama_file)
        _simpan_cache_ekstraksi_office(hash_isi, body)

    teks_lengkap = f"[Ringkasan file Excel '{nama_file}'{' -- dari cache' if dari_cache else ''}] {body}"
    return teks_lengkap, dari_cache


def _ekstrak_body_docx(isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU -- POIN 2] Parsing LOKAL file .docx pakai python-docx -- BUKAN
    dikirim mentah ke Claude API (tidak didukung). Mengambil semua
    paragraf (skip yang kosong) + isi tabel (baris digabung " | ").

    Import python-docx dilakukan LAZY di sini (bukan di top-level module)
    supaya modul ai_file_reader.py tetap bisa dipakai untuk fitur lain
    (PDF/gambar/xlsx) walau python-docx belum terinstall -- baru melempar
    error yang jelas kalau user benar-benar upload .docx.

    Return BODY saja (tanpa header) -- lihat catatan di _ekstrak_body_xlsx.
    """
    try:
        import docx  # python-docx
    except ImportError as e:
        raise RuntimeError(
            "Library 'python-docx' belum terinstall -- jalankan "
            "`pip install python-docx` untuk mendukung upload file .docx."
        ) from e

    # [BARU -- POIN 3] Sama seperti _ekstrak_body_xlsx -- cek signature
    # OLE Compound File dulu sebelum python-docx dipanggil.
    if _tampak_seperti_file_terkunci_password(isi_bytes):
        raise ValueError(
            f"'{nama_file}' terkunci password (proteksi buka-file) -- tidak bisa dibaca otomatis. "
            "Buka file ini secara manual, masukkan passwordnya, lalu simpan ulang (Save As) TANPA "
            "proteksi password sebelum upload lagi."
        )
    try:
        dokumen = docx.Document(io.BytesIO(isi_bytes))
    except Exception as e:  # noqa: BLE001
        raise _pesan_error_kemungkinan_password(nama_file, "Word", e)

    bagian: List[str] = []
    for p in dokumen.paragraphs:
        if p.text.strip():
            bagian.append(p.text)

    for i, tabel in enumerate(dokumen.tables):
        bagian.append(f"\n=== Tabel {i + 1} ===")
        for row in tabel.rows:
            bagian.append(" | ".join(cell.text.strip() for cell in row.cells))

    return "\n".join(bagian)


# ============================================================
# [BARU -- POIN 4, PROSES PARALEL] Delegasi body-extraction CPU-berat
# ke PROSES terpisah (bukan thread) -- lihat catatan benchmark di
# kirim_banyak_file_ke_ai(): openpyxl/python-docx/python-pptx murni
# Python (GIL-bound), jadi ThreadPoolExecutor SAJA tidak mempercepat
# parsing-nya sama sekali (sudah diukur langsung, bukan asumsi -- 6 file
# xlsx 8000 baris: paralel 1.42s vs sequential 1.33s, nyaris sama).
# ProcessPoolExecutor benar-benar lepas GIL karena tiap proses punya
# interpreter Python sendiri -- inilah yang dipakai di sini utk bagian
# CPU-berat SAJA (body extraction saat cache-miss), sementara cache-check
# (I/O disk, murah) tetap di thread pemanggil.
# ============================================================

def _ekstrak_body_generik(ext: str, isi_bytes: bytes, nama_file: str) -> str:
    """
    Dispatcher body-extraction TANPA cache-check, dipanggil DI DALAM
    proses worker (ProcessPoolExecutor.submit) -- HARUS didefinisikan di
    level modul (bukan closure/lambda/nested function) & semua argumen +
    nilai baliknya harus picklable (str/bytes memenuhi ini), supaya bisa
    dikirim lintas-proses lewat multiprocessing (pickle). Proses worker
    me-reimport modul ini dari awal (Python 'spawn'), jadi lazy-import
    python-docx/python-pptx/xlrd/textract di fungsi _ekstrak_body_*
    masing-masing tetap berlaku normal di proses worker.
    """
    if ext in (".xlsx", ".xlsm"):
        return _ekstrak_body_xlsx(isi_bytes, nama_file)
    if ext == ".docx":
        return _ekstrak_body_docx(isi_bytes, nama_file)
    if ext == ".pptx":
        return _ekstrak_body_pptx(isi_bytes, nama_file)
    if ext == ".xls":
        return _ekstrak_body_xls_lama(isi_bytes, nama_file)
    if ext == ".doc":
        return _ekstrak_body_doc_lama(isi_bytes, nama_file)
    raise ValueError(f"Ekstensi '{ext}' tidak didukung untuk ekstraksi proses paralel.")


# Jumlah proses worker -- default sisakan minimal 1 core utk proses utama
# (server FastAPI + thread pool I/O), bisa dioverride lewat env var kalau
# server punya banyak core & mau dorong lebih agresif. Terlalu banyak
# proses = overhead memory (tiap proses re-import pandas/openpyxl, ~50-
# 100MB per proses) tanpa manfaat sepadan kalau core fisik terbatas.
MAX_WORKER_PROSES_EKSTRAKSI = int(
    os.environ.get("AI_FILE_READER_MAX_PROSES", str(max(1, min(4, (os.cpu_count() or 2) - 1))))
)

_process_pool: Optional[concurrent.futures.ProcessPoolExecutor] = None
_process_pool_lock = threading.Lock()


def _ambil_process_pool() -> concurrent.futures.ProcessPoolExecutor:
    """
    Lazy-init PERSISTENT process pool (dibuat SEKALI, dipakai ulang lintas
    request) -- BUKAN `with ProcessPoolExecutor() as pool:` baru tiap
    panggilan, karena spawn proses baru itu sendiri mahal (~100-300ms per
    proses di Windows/spawn method) -- kalau pool dibuat-bongkar tiap
    request, overhead spawn bisa lebih besar dari waktu yang dihemat oleh
    paralelisasi itu sendiri. Double-checked locking supaya thread-safe
    kalau beberapa request FastAPI bersamaan trigger init pertama kali.

    [PENTING -- INTEGRASI FASTAPI] Init pool ini SENGAJA lazy (baru
    terjadi saat file pertama benar-benar diproses), BUKAN saat modul
    di-import -- di Windows, multiprocessing pakai method 'spawn' yang
    me-restart proses Python dari awal & re-import modul; kalau pool
    dibuat di level modul (saat import), tiap proses worker yang baru
    di-spawn akan IKUT meng-import modul ini lagi -> ikut trigger
    pembuatan pool lagi -> proses ber-spawn tak terbatas. Karena di sini
    baru dibuat saat FUNGSI dipanggil (bukan saat import), aman.
    """
    global _process_pool
    if _process_pool is None:
        with _process_pool_lock:
            if _process_pool is None:  # double-checked
                _process_pool = concurrent.futures.ProcessPoolExecutor(
                    max_workers=MAX_WORKER_PROSES_EKSTRAKSI
                )
                atexit.register(_tutup_process_pool_ekstraksi)
                logger.info(
                    f"🚀 Process pool ekstraksi file dibuat ({MAX_WORKER_PROSES_EKSTRAKSI} worker)."
                )
    return _process_pool


def _tutup_process_pool_ekstraksi() -> None:
    """[BARU] Matikan process pool dengan rapi -- didaftarkan lewat
    atexit SEBAGAI JARING PENGAMAN, tapi endpoint FastAPI (main.py)
    SEBAIKNYA memanggil ini secara EKSPLISIT di shutdown event/lifespan
    handler, supaya proses worker tidak menggantung nunggu atexit saat
    server di-restart paksa (SIGKILL, container redeploy, dst)."""
    global _process_pool
    if _process_pool is not None:
        _process_pool.shutdown(wait=False, cancel_futures=True)
        _process_pool = None


# [BARU] Timeout per file utk ekstraksi di proses terpisah -- tanpa ini,
# 1 file korup/aneh yang bikin parser hang (mis. loop tak berhenti di
# openpyxl gara-gara struktur xml rusak) akan menggantung SELAMANYA,
# menahan seluruh request (bahkan file lain yang sudah kelar tetap
# nunggu response utuh). 120 detik dipilih longgar -- file akuntansi
# wajar (puluhan ribu baris) harusnya selesai jauh di bawah itu; kalau
# sampai kena timeout, hampir pasti filenya memang bermasalah, bukan
# cuma "besar".
TIMEOUT_EKSTRAKSI_PROSES_DETIK = 120


def _jalankan_ekstraksi_di_proses(ext: str, isi_bytes: bytes, nama_file: str) -> str:
    """
    [BARU] Bungkus `_ambil_process_pool().submit(...).result()` dengan 2
    pengaman yang sebelumnya tidak ada:

    1. TIMEOUT -- kalau 1 file macet di proses worker (hang), request TIDAK
       ikut menggantung tanpa batas; dilempar TimeoutError dgn pesan jelas
       menyebut nama file yang bermasalah, supaya caller bisa laporkan ke
       user file mana yang perlu dicek ulang alih-alih server "diam saja".

    2. RECOVERY dari BrokenProcessPool -- kalau 1 proses worker crash total
       (mis. segfault dari kode C di library, kehabisan memory dipaksa
       OOM-killer), SELURUH pool jadi rusak permanen & submit berikutnya
       ke pool yang sama akan terus gagal -- tanpa recovery, server perlu
       di-restart manual utk pulih. Di sini: kalau terdeteksi rusak, pool
       lama dibuang & dibuat ulang SEKALI, lalu file yang gagal tadi
       dicoba ulang persis 1x di pool yang baru (bukan retry tak
       terbatas -- kalau file itu SENDIRI penyebab crash, retry ulang di
       pool baru cuma akan merusak pool baru lagi; makanya cuma 1x retry,
       setelah itu error asli diteruskan apa adanya).
    """
    def _submit_dan_tunggu() -> str:
        future = _ambil_process_pool().submit(_ekstrak_body_generik, ext, isi_bytes, nama_file)
        try:
            return future.result(timeout=TIMEOUT_EKSTRAKSI_PROSES_DETIK)
        except concurrent.futures.TimeoutError as e:
            raise TimeoutError(
                f"Ekstraksi '{nama_file}' melebihi batas waktu {TIMEOUT_EKSTRAKSI_PROSES_DETIK} "
                "detik -- kemungkinan file rusak/struktur tidak wajar. Coba buka & simpan ulang "
                "file ini secara manual sebelum upload lagi."
            ) from e

    try:
        return _submit_dan_tunggu()
    except concurrent.futures.process.BrokenProcessPool as e:
        logger.error(
            f"❌ Process pool ekstraksi rusak (worker crash) saat memproses '{nama_file}': {e} "
            "-- membuat ulang pool & mencoba 1x lagi."
        )
        with _process_pool_lock:
            # cek ulang di dalam lock -- kalau thread lain sudah lebih
            # dulu membuat ulang pool, jangan dibuat ulang 2x.
            global _process_pool
            if _process_pool is not None:
                _process_pool.shutdown(wait=False, cancel_futures=True)
            _process_pool = None
        return _submit_dan_tunggu()  # 1x retry di pool baru, TIDAK di-try/except lagi


def _ekstrak_teks_docx(isi_bytes: bytes, nama_file: str) -> Tuple[str, bool]:
    """
    [BARU -- POIN 2, LANJUTAN] Wrapper ber-cache di atas _ekstrak_body_docx,
    pola identik dengan _ekstrak_ringkasan_xlsx. Return (teks_lengkap, dari_cache).

    [FIX -- BUG timeout/recovery hilang di jalur satu-file] Sama seperti
    _ekstrak_ringkasan_xlsx -- lihat catatan lengkap di sana.
    """
    hash_isi = _hash_isi_file(isi_bytes)
    body = _muat_cache_ekstraksi_office(hash_isi)
    dari_cache = body is not None
    if body is None:
        body = _jalankan_ekstraksi_di_proses(".docx", isi_bytes, nama_file)
        _simpan_cache_ekstraksi_office(hash_isi, body)

    teks_lengkap = f"[Isi dokumen Word '{nama_file}'{' -- dari cache' if dari_cache else ''}]\n{body}"
    return teks_lengkap, dari_cache


def _ringkas_dataset_csv_besar(isi_bytes: bytes, nama_file: str) -> Optional[str]:
    """
    [BARU -- POIN 2] Untuk CSV berukuran besar (>MAX_KARAKTER_CSV_UNTUK_RINGKASAN),
    JANGAN dipotong mentah di karakter ke-N (versi lama) -- baris terakhir
    yang keliatan bisa saja terpotong di tengah, dan tidak ada gambaran
    STRUKTUR data (kolom apa saja, berapa baris total, ringkasan angka).

    Sebagai gantinya, diparsing pakai pandas (analog "sandbox Linux +
    pandas" yg dipakai Claude sendiri utk dataset besar) jadi ringkasan:
    jumlah baris & kolom, nama+tipe data tiap kolom, 10 baris pertama,
    dan statistik deskriptif (describe()) utk kolom numerik.

    Return None kalau parsing pandas gagal (mis. CSV formatnya aneh/rusak)
    -- caller HARUS fallback ke pemotongan teks biasa, supaya file yang
    tetap "cukup teks" walau bukan CSV standar tidak sampai gagal total.

    [FIX -- BUG NYATA] Sebelumnya `pd.read_csv(io.BytesIO(isi_bytes))`
    dipanggil TANPA deteksi delimiter/encoding -- pandas default asumsi
    koma + UTF-8. Padahal `akuntansi_ai.py::_baca_csv_sebagai_baris()`
    SUDAH punya catatan eksplisit soal ini untuk kasus yang SAMA PERSIS
    (export CSV dari Excel Indonesia): delimiter default Excel versi
    Indonesia adalah TITIK-KOMA (bukan koma, karena koma dipakai sbg
    pemisah desimal di locale ID), dan file klien sering encoding
    cp1252/latin1 (Excel Windows Indonesia), bukan UTF-8 murni.

    Sebelum fix ini, CSV titik-koma TIDAK melempar error sama sekali --
    pandas "berhasil" membaca tapi menganggap SATU BARIS UTUH sebagai
    SATU KOLOM (delimiter tidak cocok, tapi parsing tetap "sukses" secara
    teknis). Hasilnya: ringkasan yang dikirim ke AI kelihatan valid
    (`df.shape` dsb ada) tapi sebenarnya rusak total -- tidak ada
    exception yang bisa ditangkap try/except, jadi bug ini AKAN LOLOS
    tanpa pernah muncul di log mana pun. Fix: sniff delimiter + fallback
    encoding dulu (pola identik dgn akuntansi_ai.py), baru serahkan hasil
    decode+delimiter yang benar ke pandas.
    """
    import csv as _csv

    teks = None
    for enc in ("utf-8-sig", "cp1252", "latin1"):
        try:
            teks = isi_bytes.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if teks is None:
        teks = isi_bytes.decode("utf-8", errors="replace")

    sample = "\n".join(teks.splitlines()[:5])
    try:
        delimiter = _csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
    except _csv.Error:
        delimiter = ";" if sample.count(";") > sample.count(",") else ","

    try:
        df = pd.read_csv(io.StringIO(teks), sep=delimiter)
    except Exception as e:  # noqa: BLE001
        logger.warning(f"Gagal ringkas '{nama_file}' sebagai CSV lewat pandas ({e}) -- fallback ke potong teks biasa.")
        return None

    bagian: List[str] = [
        f"[Ringkasan dataset '{nama_file}' -- {df.shape[0]} baris x {df.shape[1]} kolom]",
        "\nKolom & tipe data:",
        df.dtypes.astype(str).to_string(),
        "\n10 baris pertama:",
        df.head(10).to_string(),
    ]
    kolom_numerik = df.select_dtypes(include="number")
    if not kolom_numerik.empty:
        bagian.append("\nStatistik kolom numerik:")
        bagian.append(kolom_numerik.describe().to_string())

    return "\n".join(bagian)


def _ringkas_dataset_json_besar(isi_bytes: bytes, nama_file: str) -> Optional[str]:
    """
    [BARU -- POIN 2 TAMBAHAN] Analog `_ringkas_dataset_csv_besar` di atas,
    tapi untuk .json besar. Dua kasus dibedakan:

    1. JSON berbentuk LIST OF OBJECTS (array data tabular, mis. hasil
       export data ribuan baris) -- diparsing pakai `json.load` lalu
       `pandas.json_normalize` (meratakan nested dict jadi kolom
       "a.b.c"), baru diringkas PERSIS seperti CSV: shape, kolom+tipe,
       10 baris pertama, describe() kolom numerik.

    2. JSON strukturnya BUKAN tabular (dict nested/kompleks, atau list
       berisi tipe campuran/bukan dict) -- fallback ke ringkasan
       kunci-kunci teratas saja (nama kunci top-level + tipe nilainya,
       dibatasi MAX_KUNCI_JSON_DITAMPILKAN kunci) supaya AI setidaknya
       tahu bentuk/skema datanya walau tidak semua isinya ditampilkan.

    Return None kalau JSON tidak valid sama sekali (parse gagal) --
    caller HARUS fallback ke _potong_teks_jika_perlu biasa (sama pola
    seperti _ringkas_dataset_csv_besar).
    """
    import json as _json

    try:
        data = _json.loads(isi_bytes.decode("utf-8", errors="replace"))
    except _json.JSONDecodeError as e:
        logger.warning(f"Gagal ringkas '{nama_file}' sebagai JSON ({e}) -- fallback ke potong teks biasa.")
        return None

    # --- Kasus 1: list of objects (tabular) ---
    if isinstance(data, list) and data and all(isinstance(item, dict) for item in data):
        try:
            df = pd.json_normalize(data)
        except Exception as e:  # noqa: BLE001 -- struktur nested tak terduga, jangan gagal total
            logger.warning(f"Gagal json_normalize '{nama_file}' ({e}) -- fallback ke ringkasan kunci saja.")
        else:
            bagian: List[str] = [
                f"[Ringkasan dataset JSON '{nama_file}' -- {df.shape[0]} baris x {df.shape[1]} kolom "
                "(array of objects, nested key diratakan jadi 'induk.anak')]",
                "\nKolom & tipe data:",
                df.dtypes.astype(str).to_string(),
                "\n10 baris pertama:",
                df.head(10).to_string(),
            ]
            kolom_numerik = df.select_dtypes(include="number")
            if not kolom_numerik.empty:
                bagian.append("\nStatistik kolom numerik:")
                bagian.append(kolom_numerik.describe().to_string())
            return "\n".join(bagian)

    # --- Kasus 2: bukan tabular (dict nested/kompleks, atau list campuran) ---
    return _ringkas_struktur_json_non_tabular(data, nama_file)


def _ringkas_struktur_json_non_tabular(data: Any, nama_file: str) -> str:
    """Ringkasan kunci-kunci teratas untuk JSON yang bukan array-of-objects
    tabular -- dict besar (mis. konfigurasi/response API nested) atau list
    berisi tipe campuran. Tidak mendalami isi nested lebih dari 1 level
    (cukup nama kunci + tipe nilainya) -- ini SENGAJA ringkasan skema,
    bukan dump seluruh isi, supaya tetap ringkas untuk struktur yang bisa
    saja dalam sekali."""
    bagian: List[str] = [f"[Ringkasan struktur JSON '{nama_file}' -- bukan tabular, ditampilkan skema saja]"]

    if isinstance(data, dict):
        kunci = list(data.keys())
        bagian.append(f"Objek dengan {len(kunci)} kunci top-level:")
        for k in kunci[:MAX_KUNCI_JSON_DITAMPILKAN]:
            v = data[k]
            tipe = type(v).__name__
            info_tambahan = f", panjang={len(v)}" if isinstance(v, (list, dict, str)) else ""
            bagian.append(f"  - {k}: {tipe}{info_tambahan}")
        if len(kunci) > MAX_KUNCI_JSON_DITAMPILKAN:
            bagian.append(f"  ... ({len(kunci) - MAX_KUNCI_JSON_DITAMPILKAN} kunci lain tidak ditampilkan)")
    elif isinstance(data, list):
        bagian.append(f"Array dengan {len(data)} item (tipe tidak seragam/bukan seluruhnya objek).")
        if data:
            bagian.append(f"Contoh item pertama (tipe {type(data[0]).__name__}): {str(data[0])[:500]}")
    else:
        bagian.append(f"Nilai tunggal (tipe {type(data).__name__}): {str(data)[:500]}")

    return "\n".join(bagian)


def _potong_teks_jika_perlu(isi_teks: str) -> Tuple[str, bool]:
    """[FIX -- GAP 3: FILE TEKS TIDAK DIBATASI] Potong teks yang terlalu
    panjang SEBELUM dikirim ke API -- versi dasar sebelumnya cuma
    berkomentar "pertimbangkan potong dulu" tanpa benar-benar melakukannya,
    jadi file teks besar (ribuan baris) bisa mendorong context window
    penuh sendirian atau bikin biaya token melonjak tanpa peringatan.
    MAX_KARAKTER_TEKS adalah batas kasar (~100rb token) -- cukup longgar
    untuk kebanyakan file konfigurasi/data teks, tapi bukan tak terbatas.
    Return (teks_final, dipotong_atau_tidak)."""
    if len(isi_teks) <= MAX_KARAKTER_TEKS:
        return isi_teks, False
    return isi_teks[:MAX_KARAKTER_TEKS], True


def _siapkan_teks_dari_bytes_teks(isi_bytes: bytes, nama_file: str) -> Tuple[str, bool]:
    """
    [BARU -- POIN 2] Dipakai untuk tipe "teks" (md/txt/csv/json/html).
    KHUSUS untuk .csv yang besar (>MAX_KARAKTER_CSV_UNTUK_RINGKASAN),
    coba diringkas pakai pandas dulu (lihat _ringkas_dataset_csv_besar) --
    "dataset besar" diperlakukan beda dari file teks kecil biasa, BUKAN
    dipotong mentah di karakter ke-N. Kalau parsing pandas gagal (None),
    fallback ke potong teks biasa seperti file lain.

    Return (teks_final, sudah_diringkas_atau_dipotong).
    """
    ext = Path(nama_file).suffix.lower()
    if ext == ".csv" and len(isi_bytes) > MAX_KARAKTER_CSV_UNTUK_RINGKASAN:
        ringkasan = _ringkas_dataset_csv_besar(isi_bytes, nama_file)
        if ringkasan is not None:
            return ringkasan, True

    # [BARU -- POIN 2 TAMBAHAN] .json besar diringkas (tabular via pandas
    # ATAU skema kunci top-level), BUKAN dipotong mentah di karakter ke-N
    # seperti sebelumnya -- lihat _ringkas_dataset_json_besar.
    if ext == ".json" and len(isi_bytes) > MAX_KARAKTER_JSON_UNTUK_RINGKASAN:
        ringkasan = _ringkas_dataset_json_besar(isi_bytes, nama_file)
        if ringkasan is not None:
            return ringkasan, True

    isi_teks = isi_bytes.decode("utf-8", errors="replace")
    return _potong_teks_jika_perlu(isi_teks)


def siapkan_konten_pesan_dari_file(isi_bytes: bytes, nama_file: str, pertanyaan: str) -> Any:
    """
    [BARU] Logic deteksi-tipe + validasi + ekstraksi/base64-encode, DIPISAH
    dari kirim_file_ke_ai supaya bisa dipakai ULANG oleh versi non-stream
    (kirim_file_ke_ai) MAUPUN versi stream (kirim_file_ke_ai_stream) tanpa
    duplikasi 5 cabang if/elif yang sama di 2 tempat (rawan salah satu
    diedit belakangan tapi yang lain lupa ikut diubah).

    SENGAJA public (tanpa underscore) -- endpoint FastAPI streaming di
    main.py PERLU memanggil ini SENDIRI, SEBELUM membuka StreamingResponse:
    kalau file tidak valid (ValueError) atau API key belum diset
    (RuntimeError), errornya harus ketahuan & jadi HTTPException 400/500
    yang BENAR *sebelum* header response mulai dikirim ke klien -- begitu
    StreamingResponse mulai jalan, status code 200 sudah terkirim duluan,
    tidak bisa diganti lagi jadi 400 di tengah jalan kalau ternyata gagal.

    Return `content` siap pakai untuk `messages=[{"role": "user",
    "content": <ini>}]` -- string polos (file teks/office/office_lama)
    ATAU list content block (gambar/PDF, base64 + teks pertanyaan).

    Melempar ValueError/RuntimeError sama seperti kirim_file_ke_ai.
    """
    tipe = deteksi_tipe_file(nama_file)

    if tipe == "teks":
        isi_teks, dipotong_atau_ringkas = _siapkan_teks_dari_bytes_teks(isi_bytes, nama_file)
        prompt = f"Berikut isi file '{nama_file}':\n\n{isi_teks}\n\n{pertanyaan}"
        if dipotong_atau_ringkas:
            prompt += "\n\n(Catatan: isi file di atas sudah diringkas/dipotong karena terlalu besar -- lihat bagian yang ditampilkan saja.)"
        return prompt

    if tipe == "office":
        _validasi_office(isi_bytes, nama_file)
        ext = Path(nama_file).suffix.lower()
        if ext in (".xlsx", ".xlsm"):
            isi_ekstraksi, _dari_cache = _ekstrak_ringkasan_xlsx(isi_bytes, nama_file)
        elif ext == ".docx":
            isi_ekstraksi, _dari_cache = _ekstrak_teks_docx(isi_bytes, nama_file)
        else:  # .pptx
            isi_ekstraksi, _dari_cache = _ekstrak_teks_pptx(isi_bytes, nama_file)
        isi_ekstraksi, dipotong = _potong_teks_jika_perlu(isi_ekstraksi)
        prompt = f"{isi_ekstraksi}\n\n{pertanyaan}"
        if dipotong:
            prompt += "\n\n(Catatan: hasil ekstraksi dipotong karena terlalu panjang.)"
        return prompt

    if tipe == "office_lama":
        _validasi_office(isi_bytes, nama_file)
        isi_ekstraksi, _dari_cache = ekstrak_office_lama(isi_bytes, nama_file)
        isi_ekstraksi, dipotong = _potong_teks_jika_perlu(isi_ekstraksi)
        prompt = f"{isi_ekstraksi}\n\n{pertanyaan}"
        if dipotong:
            prompt += "\n\n(Catatan: hasil ekstraksi dipotong karena terlalu panjang.)"
        return prompt

    if tipe == "gambar":
        _validasi_gambar(isi_bytes, nama_file)
        media_type = _MEDIA_TYPE_GAMBAR[Path(nama_file).suffix.lower()]
        base64_data = base64.b64encode(isi_bytes).decode("utf-8")
        return [
            {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": base64_data}},
            {"type": "text", "text": pertanyaan},
        ]

    if tipe == "pdf":
        # [FIX -- SEMENTARA, lihat blok "_ekstrak_teks_pdf_penuh" di atas]
        # SEBELUMNYA PDF selalu dikirim mentah sbg base64 content block
        # "document" (baris di-comment di bawah) -- itu HANYA bisa dibaca
        # Claude. Sekarang, selama ANTHROPIC_API_KEY belum ada saldo,
        # PDF diekstrak jadi teks UTUH dulu supaya bisa juga dijawab
        # provider OpenAI-compatible (Groq) yg cuma terima teks polos.
        # Ukuran file TETAP divalidasi (_validasi_pdf) walau tidak lagi
        # dikirim sbg base64 -- supaya tidak coba proses file yg sudah
        # jelas melebihi batas wajar.
        _validasi_pdf(isi_bytes, nama_file)
        isi_teks = _ekstrak_teks_pdf_penuh(isi_bytes, nama_file)
        isi_teks, dipotong = _potong_teks_jika_perlu(isi_teks)
        prompt = f"Berikut isi PDF '{nama_file}' (hasil ekstraksi teks):\n\n{isi_teks}\n\n{pertanyaan}"
        if dipotong:
            prompt += "\n\n(Catatan: isi PDF di atas dipotong karena terlalu panjang.)"
        return prompt
        # --- Jalur LAMA (kirim PDF mentah, HANYA jalan dgn Claude) ---
        # base64_data = _encode_pdf_base64_dengan_cache(isi_bytes)
        # return [
        #     {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": base64_data}},
        #     {"type": "text", "text": pertanyaan},
        # ]

    raise ValueError(
        f"Tipe file '{Path(nama_file).suffix}' tidak didukung. "
        f"Didukung: {sorted(_EKSTENSI_TEKS | _EKSTENSI_GAMBAR | _EKSTENSI_OFFICE | _EKSTENSI_OFFICE_LAMA | {'.pdf'})}"
    )


def tutup_process_pool_ekstraksi() -> None:
    """
    [BARU -- POIN 4] Titik masuk PUBLIK (tanpa underscore) untuk dipanggil
    dari shutdown event / lifespan handler FastAPI di main.py, mis.:

        from modules.ai_file_reader import tutup_process_pool_ekstraksi

        @app.on_event("shutdown")          # atau lifespan context manager
        def _shutdown():
            tutup_process_pool_ekstraksi()

    Tanpa ini, process pool tetap akan dimatikan via atexit sbg jaring
    pengaman (lihat _ambil_process_pool), tapi shutdown eksplisit lebih
    bersih & lebih cepat pada graceful shutdown server (uvicorn --reload,
    container redeploy, dst) -- proses worker langsung diberi sinyal
    berhenti alih-alih menunggu interpreter utama benar-benar exit.
    """
    _tutup_process_pool_ekstraksi()


def kirim_file_ke_ai(
    isi_bytes: bytes, nama_file: str, pertanyaan: str, model: str = MODEL_DEFAULT,
) -> str:
    """
    Titik masuk utama untuk dipakai dari endpoint FastAPI -- terima BYTES
    (hasil `await file.read()`) + nama file asli, auto-deteksi cara
    mengirimnya ke Claude:
    - teks kecil: teks polos di prompt
    - CSV besar ("dataset besar"): diringkas pakai pandas dulu (POIN 2)
    - gambar/PDF: base64 content block, dikirim mentah
    - xlsx/xlsm/docx/pptx & xls/doc lama ("file besar/biner", POIN 2):
      diparsing LOKAL pakai openpyxl/python-docx/python-pptx/xlrd/textract
      dulu, hasil ekstraksinya (teks) yang dikirim, BUKAN file binernya --
      Claude API tidak menerima format Office mentah.

    Melempar ValueError kalau tipe file tidak didukung ATAU melanggar
    limit (resmi Claude API utk PDF/gambar, atau limit kita sendiri utk
    Office) -- caller (endpoint) yang menerjemahkan ValueError ini
    jadi HTTPException 400.

    Untuk versi yang jawabannya di-STREAM bertahap (bukan tunggu selesai
    semua), lihat kirim_file_ke_ai_stream() di bawah.

    [FIX -- SEMENTARA] Selama ANTHROPIC_API_KEY belum ada saldo: SEMUA
    tipe file yang content-nya berupa TEKS POLOS (teks/csv/json/html, PDF
    hasil ekstraksi, Office hasil ekstraksi -- lihat
    siapkan_konten_pesan_dari_file) dikirim lewat _tanya_teks_ke_provider()
    yang coba Claude dulu, fallback Groq kalau kosong/gagal. HANYA tipe
    "gambar" yang TETAP wajib Claude (content-nya list berisi image
    content block -- Groq model teks di sini tidak punya kemampuan vision,
    tidak bisa "melihat" gambar sama sekali).
    """
    content = siapkan_konten_pesan_dari_file(isi_bytes, nama_file, pertanyaan)

    if isinstance(content, str):
        return _tanya_teks_ke_provider(content)

    # content berupa list -> ada content block gambar, WAJIB Claude (vision)
    response = _panggil_dengan_retry(
        model=model,
        max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
        system=SYSTEM_PROMPT_FILE_READER,
        messages=[{"role": "user", "content": content}],
    )
    return response.content[0].text


def kirim_file_ke_ai_stream(content: Any, model: str = MODEL_DEFAULT):
    """
    [BARU] Versi STREAMING -- generator yang yield potongan teks jawaban
    SEGERA saat diterima dari provider, bukan menunggu jawaban lengkap
    baru dikembalikan sekaligus (mirip pengalaman "mengetik" bertahap yang
    dipakai Claude AI sendiri) -- berguna terutama utk pertanyaan panjang/
    kompleks di atas hasil ekstraksi file besar.

    [PENTING] Terima `content` yang SUDAH SIAP (hasil
    siapkan_konten_pesan_dari_file(), BUKAN raw bytes) -- SENGAJA dipisah
    dari tahap ekstraksi/validasi supaya caller (endpoint FastAPI) bisa
    panggil siapkan_konten_pesan_dari_file() dulu SEBELUM membuka
    StreamingResponse, dan menangani ValueError/RuntimeError di sana jadi
    HTTPException 400/500 yang benar. Kalau ekstraksi/validasi dilakukan
    DI DALAM generator ini, errornya baru ketahuan setelah StreamingResponse
    sudah mulai jalan (status 200 sudah terkirim ke klien) -- tidak bisa
    diganti jadi 400 lagi di tengah jalan.

    [FIX -- SEMENTARA] Routing ditentukan dari TIPE content, bukan dari
    nama file: kalau `content` berupa string (teks/csv/json/html, PDF
    hasil ekstraksi, Office hasil ekstraksi), jalur ini IKUT fallback ke
    Groq (lewat _konfigurasi_provider_file_reader) kalau ANTHROPIC_API_KEY
    kosong/tidak ada saldo. Kalau `content` berupa list (ada content block
    gambar), TETAP wajib Claude saja -- Groq model teks di sini tidak
    punya kemampuan vision.

    CATATAN RETRY: retry manual (_panggil_dengan_retry) SENGAJA tidak
    dipakai untuk jalur Claude di sini -- begitu stream mulai jalan &
    sebagian teks sudah terkirim ke caller, tidak ada cara "retry dari
    awal" yang bersih tanpa caller menerima teks dobel. Untuk fallback
    Groq, pindah provider HANYA dilakukan kalau BELUM ADA token sama
    sekali yang keluar dari provider sebelumnya (sama pola dgn
    tanya_ai_stream di akuntansi_ai.py) -- begitu ada token yang sudah
    terkirim, error berikutnya dilempar apa adanya, tidak pindah provider
    lagi (mencegah jawaban dobel/tercampur ke caller).

    Yields: potongan teks (str) satu per satu, sesuai urutan diterima.
    """
    if not isinstance(content, str):
        # content berupa list -> ada content block gambar, WAJIB Claude (vision)
        client = _ambil_client()
        with client.messages.stream(
            model=model,
            max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
            system=SYSTEM_PROMPT_FILE_READER,
            messages=[{"role": "user", "content": content}],
        ) as stream:
            for teks in stream.text_stream:
                yield teks
        return

    # --- content berupa teks polos: coba tiap provider di
    # _konfigurasi_provider_file_reader, pindah ke provider berikutnya
    # HANYA kalau belum ada token yang keluar.
    daftar_provider = _konfigurasi_provider_file_reader()
    if not daftar_provider:
        raise RuntimeError(
            "Tidak ada API key aktif untuk membaca file (ANTHROPIC_API_KEY / "
            "GROQ_API_KEY_KATEGORISASI / GROQ_API_KEY semuanya kosong)."
        )

    error_terakhir: Optional[Exception] = None
    for konfig in daftar_provider:
        sudah_ada_token = False
        try:
            if konfig["tipe"] == "anthropic":
                client = _ambil_client()
                with client.messages.stream(
                    model=konfig["model"],
                    max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
                    system=SYSTEM_PROMPT_FILE_READER,
                    messages=[{"role": "user", "content": content}],
                ) as stream:
                    for teks in stream.text_stream:
                        sudah_ada_token = True
                        yield teks
                return
            else:  # "openai_compatible" -- Groq (system prompt sbg pesan
                # ber-role "system" di messages, bukan parameter terpisah)
                import openai
                client = openai.OpenAI(api_key=konfig["api_key"], base_url=konfig["base_url"], max_retries=0)
                stream = client.chat.completions.create(
                    model=konfig["model"],
                    messages=[
                        {"role": "system", "content": SYSTEM_PROMPT_FILE_READER},
                        {"role": "user", "content": content},
                    ],
                    stream=True,
                    **konfig.get("extra_params", {}),
                )
                for potongan in stream:
                    delta = potongan.choices[0].delta.content
                    if delta:
                        sudah_ada_token = True
                        yield delta
                return
        except Exception as e:  # noqa: BLE001
            error_terakhir = e
            if sudah_ada_token:
                raise
            logger.warning(f"⚠️ Provider '{konfig['nama']}' gagal baca file (stream): {e} -- coba provider berikutnya.")
            continue

    raise error_terakhir or RuntimeError("Gagal memanggil semua provider file reader (stream).")


# ============================================================
# E. [BARU -- GAP 5: SATU FILE PER REQUEST] Kirim BANYAK file sekaligus
# dalam SATU pertanyaan -- mis. "bandingkan isi file A dan file B" atau
# "gabungkan semua rekening koran ini jadi 1 rekap" -- API mendukung
# banyak content block gambar/dokumen dalam 1 pesan, versi dasar
# sebelumnya cuma bisa 1 file per panggilan sehingga tidak bisa dipakai
# untuk kasus lintas-file seperti itu.
# ============================================================

def _olah_satu_file_untuk_multi(
    isi_bytes: bytes, nama_file: str, tipe: str,
) -> Tuple[str, Any, int]:
    """
    [BARU -- POIN 4] Kerja SATU file (validasi + ekstraksi/encode) --
    diekstrak jadi fungsi sendiri supaya bisa dijalankan PARALEL per file
    lewat ThreadPoolExecutor di kirim_banyak_file_ke_ai(), bukan bergiliran
    satu-satu. Ini yang sebenarnya memberi percepatan nyata: kalau user
    upload 5 xlsx besar tanpa cache-hit, openpyxl/pandas mem-parsing
    kelimanya BERSAMAAN (thread lain jalan selagi 1 thread menunggu I/O
    disk / lepas GIL di kode C pandas/openpyxl), bukan menunggu file 1
    selesai baru mulai file 2.

    Return (kind, payload, ukuran_bytes_biner):
      - kind "teks": payload = str blok teks siap gabung ke prompt
      - kind "gambar" / "pdf": payload = content block dict siap kirim,
        ukuran_bytes_biner = len(isi_bytes) (dipakai caller utk cek
        limit 32MB total SETELAH semua thread selesai)
    File teks/office tidak masuk hitungan ukuran_bytes_biner (0) karena
    tidak dikirim sbg base64 content block.

    Melempar ValueError/RuntimeError apa adanya -- caller
    (kirim_banyak_file_ke_ai) yang menangkap & meneruskan ke pemanggil,
    persis perilaku sebelumnya waktu masih sequential.
    """
    if tipe == "teks":
        isi_teks, _ = _siapkan_teks_dari_bytes_teks(isi_bytes, nama_file)
        return "teks", f"=== Isi file '{nama_file}' ===\n{isi_teks}", 0

    if tipe == "office":
        # [BARU -- POIN 2, LANJUTAN] xlsx/docx/pptx diparsing lokal --
        # hasil teksnya digabung ke blok teks di prompt, bukan jadi
        # content block base64.
        _validasi_office(isi_bytes, nama_file)
        ext = Path(nama_file).suffix.lower()
        hash_isi = _hash_isi_file(isi_bytes)
        body = _muat_cache_ekstraksi_office(hash_isi)
        dari_cache = body is not None
        if body is None:
            # [BARU -- POIN 4, PROSES PARALEL] cache-miss -> delegasikan
            # body extraction (kerja CPU berat) ke PROSES terpisah, bukan
            # diparsing inline di thread ini -- lihat catatan benchmark
            # di _ekstrak_body_generik/kirim_banyak_file_ke_ai kenapa ini
            # perlu proses, bukan cukup thread. Dibungkus timeout +
            # recovery BrokenProcessPool -- lihat _jalankan_ekstraksi_di_proses.
            body = _jalankan_ekstraksi_di_proses(ext, isi_bytes, nama_file)
            _simpan_cache_ekstraksi_office(hash_isi, body)
        cache_note = " -- dari cache" if dari_cache else ""
        if ext in (".xlsx", ".xlsm"):
            isi_ekstraksi = f"[Ringkasan file Excel '{nama_file}'{cache_note}] {body}"
        elif ext == ".docx":
            isi_ekstraksi = f"[Isi dokumen Word '{nama_file}'{cache_note}]\n{body}"
        else:  # .pptx
            isi_ekstraksi = f"[Isi presentasi PowerPoint '{nama_file}'{cache_note}] {body}"
        isi_ekstraksi, _ = _potong_teks_jika_perlu(isi_ekstraksi)
        return "teks", isi_ekstraksi, 0

    if tipe == "office_lama":
        _validasi_office(isi_bytes, nama_file)
        ext = Path(nama_file).suffix.lower()
        hash_isi = _hash_isi_file(isi_bytes)
        body = _muat_cache_ekstraksi_office(hash_isi)
        dari_cache = body is not None
        if body is None:
            # [BARU -- POIN 4, PROSES PARALEL] sama pola dgn cabang
            # "office" di atas -- .xls lewat xlrd (murni Python, GIL-
            # bound juga) tetap diuntungkan proses terpisah; .doc lewat
            # textract (subprocess antiword) tetap aman didelegasikan.
            # Timeout + recovery sama, lihat _jalankan_ekstraksi_di_proses.
            body = _jalankan_ekstraksi_di_proses(ext, isi_bytes, nama_file)
            _simpan_cache_ekstraksi_office(hash_isi, body)
        label = "Excel" if ext == ".xls" else "Word"
        cache_note = " -- dari cache" if dari_cache else ""
        isi_ekstraksi = f"[Isi dokumen {label} lama '{nama_file}'{cache_note}]\n{body}"
        isi_ekstraksi, _ = _potong_teks_jika_perlu(isi_ekstraksi)
        return "teks", isi_ekstraksi, 0

    if tipe == "gambar":
        _validasi_gambar(isi_bytes, nama_file)
        media_type = _MEDIA_TYPE_GAMBAR[Path(nama_file).suffix.lower()]
        blok = {
            "type": "image",
            "source": {"type": "base64", "media_type": media_type, "data": base64.b64encode(isi_bytes).decode("utf-8")},
        }
        return "gambar", blok, len(isi_bytes)

    if tipe == "pdf":
        _validasi_pdf(isi_bytes, nama_file)
        blok = {
            "type": "document",
            "source": {"type": "base64", "media_type": "application/pdf", "data": _encode_pdf_base64_dengan_cache(isi_bytes)},
        }
        return "pdf", blok, len(isi_bytes)

    raise ValueError(
        f"Tipe file '{Path(nama_file).suffix}' (dari '{nama_file}') tidak didukung. "
        f"Didukung: {sorted(_EKSTENSI_TEKS | _EKSTENSI_GAMBAR | _EKSTENSI_OFFICE | _EKSTENSI_OFFICE_LAMA | {'.pdf'})}"
    )


# [BARU -- POIN 4] Batas jumlah thread paralel utk ekstraksi multi-file --
# tidak sekaligus tak terbatas (mis. 50 file dalam 1 request tidak boleh
# buka 50 thread openpyxl/pandas sekaligus, bisa bikin CPU/memory server
# thrashing) -- 8 dipilih sbg titik wajar: cukup utk paralelkan kasus
# umum (2-10 file dibandingkan/digabung sekaligus) tanpa oversubscribe.
MAX_WORKER_PARALEL_MULTI_FILE = 8


def kirim_banyak_file_ke_ai(
    daftar_file: List[Tuple[bytes, str]],
    pertanyaan: str,
    model: str = MODEL_DEFAULT,
) -> str:
    """
    Args:
        daftar_file: list of (isi_bytes, nama_file) -- boleh campur
            teks/gambar/PDF dalam 1 panggilan.
        pertanyaan: pertanyaan yang berlaku utk SEMUA file di atas
            (mis. "bandingkan file A dan B", "rekap semua transaksi
            dari file-file berikut").
        model: default MODEL_DEFAULT.

    Melempar ValueError kalau: ada file tidak didukung, ada file yang
    melanggar limit ukuran individual, jumlah gambar > MAX_GAMBAR_PER_REQUEST,
    atau total ukuran gabungan (PDF+gambar, base64-encoded) > 32MB (limit
    request Claude API dihitung PER REQUEST, bukan per file, kalau lebih
    dari 1 dokumen/gambar dikirim sekaligus).

    File teks digabung jadi satu blok teks di awal prompt (dipisah nama
    file), gambar & PDF masing-masing jadi content block base64 sendiri.

    [BARU -- POIN 4] Validasi/ekstraksi PER FILE dijalankan PARALEL
    (ThreadPoolExecutor, lihat _olah_satu_file_untuk_multi) -- bukan
    bergiliran satu-satu seperti sebelumnya. Urutan hasil di prompt akhir
    TETAP mengikuti urutan `daftar_file` (pakai executor.map, yang
    menjaga urutan output = urutan input walau thread selesai tidak
    berurutan), jadi perilaku yang terlihat pemanggil sama persis
    dengan versi sequential -- cuma lebih cepat.
    """
    if not daftar_file:
        raise ValueError("Tidak ada file untuk dikirim.")

    # --- Pre-pass murah (bukan kerja berat, cuma baca ekstensi & len(bytes),
    # BUKAN encode base64) -- cek limit jumlah gambar & PERKIRAAN ukuran
    # payload gabungan SEBELUM buka thread pool, supaya gagal cepat tanpa
    # buang kerja paralel (termasuk encode base64 yang tidak murah untuk
    # PDF besar) untuk file-file lain kalau sudah pasti akan ditolak.
    #
    # [FIX -- POIN 4] Sebelumnya pre-pass ini CUMA cek jumlah gambar --
    # cek ukuran 32MB baru dilakukan SETELAH executor.map selesai (yaitu
    # setelah SEMUA file, termasuk PDF besar, sudah selesai di-encode
    # base64 di masing-masing thread) meski komentar lama mengklaim
    # "supaya gagal SEBELUM proses encode selesai". Sekarang dicek di
    # SINI dulu (byte mentah, sebelum encode) -- kalau sudah pasti > 32MB
    # setelah estimasi overhead base64, gagal cepat TANPA membuang CPU
    # meng-encode file yang toh bakal ditolak.
    tipe_per_file = [deteksi_tipe_file(nama) for _, nama in daftar_file]
    jumlah_gambar = sum(1 for t in tipe_per_file if t == "gambar")
    if jumlah_gambar > MAX_GAMBAR_PER_REQUEST:
        raise ValueError(
            f"Terlalu banyak gambar dalam 1 request ({jumlah_gambar}), "
            f"limit Claude API {MAX_GAMBAR_PER_REQUEST} gambar per request. "
            "Kirim dalam beberapa panggilan terpisah."
        )

    total_bytes_biner_mentah = sum(
        len(isi) for (isi, _nama), tipe in zip(daftar_file, tipe_per_file)
        if tipe in ("gambar", "pdf")
    )
    perkiraan_payload_awal = int(total_bytes_biner_mentah * 1.34)
    if perkiraan_payload_awal > MAX_PDF_BYTES:
        raise ValueError(
            f"Total ukuran file gabungan (PDF/gambar) diperkirakan "
            f"{perkiraan_payload_awal / (1024*1024):.1f} MB setelah encoding, melebihi "
            f"limit request Claude API {MAX_PDF_BYTES // (1024*1024)} MB. "
            "Kirim file dalam beberapa panggilan lebih kecil."
        )

    import concurrent.futures

    n_worker = min(MAX_WORKER_PARALEL_MULTI_FILE, len(daftar_file))
    with concurrent.futures.ThreadPoolExecutor(max_workers=n_worker) as executor:
        hasil_per_file = list(executor.map(
            lambda args: _olah_satu_file_untuk_multi(args[0][0], args[0][1], args[1]),
            zip(daftar_file, tipe_per_file),
        ))

    blok_teks: List[str] = []
    content_blocks: List[Dict[str, Any]] = []
    total_bytes_biner = 0  # perkiraan payload PDF+gambar sebelum base64

    for kind, payload, ukuran in hasil_per_file:
        if kind == "teks":
            blok_teks.append(payload)
        else:  # "gambar" atau "pdf"
            content_blocks.append(payload)
            total_bytes_biner += ukuran

    # [FIX -- GAP 2, versi multi-file] JARING PENGAMAN KEDUA -- gerbang
    # UTAMA sudah dicek lebih awal (lihat perkiraan_payload_awal SEBELUM
    # thread pool dibuka), ini cuma verifikasi ulang pakai ukuran base64
    # AKTUAL (bukan perkiraan dari byte mentah) setelah semua file selesai
    # diproses, jaga-jaga kalau ada selisih pembulatan antara estimasi
    # awal & hasil encode sungguhan.
    perkiraan_payload = int(total_bytes_biner * 1.34)
    if perkiraan_payload > MAX_PDF_BYTES:
        raise ValueError(
            f"Total ukuran file gabungan (PDF/gambar) diperkirakan "
            f"{perkiraan_payload / (1024*1024):.1f} MB setelah encoding, melebihi "
            f"limit request Claude API {MAX_PDF_BYTES // (1024*1024)} MB. "
            "Kirim file dalam beberapa panggilan lebih kecil."
        )

    prompt_teks = pertanyaan
    if blok_teks:
        prompt_teks = "\n\n".join(blok_teks) + f"\n\n{pertanyaan}"

    content_blocks.append({"type": "text", "text": prompt_teks})

    response = _panggil_dengan_retry(
        model=model,
        max_tokens=MAX_TOKENS_JAWABAN_DEFAULT,
        system=SYSTEM_PROMPT_FILE_READER,
        messages=[{"role": "user", "content": content_blocks}],
    )
    return response.content[0].text


# ============================================================
# Contoh pemakaian langsung (uji coba dari command line)
# ============================================================
if __name__ == "__main__":
    # python modules/ai_file_reader.py
    hasil = kirim_pdf_ke_ai(
        "contoh_rekening_koran.pdf",
        "Ekstrak semua transaksi jadi tabel: tanggal, keterangan, mutasi debet, mutasi kredit."
    )
    print(hasil)